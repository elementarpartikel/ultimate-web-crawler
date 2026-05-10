#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Webbdammsugare Pro (v7.0)
Skapad av Fredrik Eriksson

Asynkron webbcrawler med GUI (CustomTkinter) och CLI-/serverläge.

Funktioner:
  - Hybrid-motor: aiohttp + Playwright-fallback för JavaScript-tunga sidor
  - Inkrementell crawl med conditional GET (ETag / If-Modified-Since)
  - Riktig CookieJar med stöd för SAML/SSO-inloggning via Playwright
  - Per-domän rate limiter och prioritetskö med boostning/penalty
  - Sitemap-parser (XML, gzip, sitemap-index)
  - robots.txt-stöd med Crawl-Delay
  - Login-detektor (URL-redirect, HTTP-status, innehållsheuristik)
  - Dokumentnedladdning (PDF, Word, Excel, PPTX, ZIP m.fl.)
  - Dokument → Markdown-konvertering med källinformation i brödtexten
  - Dokument-manifest (manifest.json) som kopplar filer till ursprungssida
  - CMS-boilerplate-rensning (Sitevision m.fl.)
  - PII-tvätt (e-post, personnummer, telefon, IP)
  - Semantisk chunkning med URL per chunk
  - Trafilatura-integration för bättre textextraktion
  - Sitevision-anpassad URL-normalisering (sv.*, state, logout)
  - Smart titelextraktion (link_text → metadata → filnamn → referer_title)
  - Generisk-länktext-filter ("Ladda ner fil", "Download" m.fl.)
  - Batched DB-commits, valbar samtidighet, snabb stopphantering
"""

import argparse
import asyncio
import csv
import gzip
import hashlib
import json
import logging
import os
import posixpath
import queue
import re
import sqlite3
import sys
import threading
import time
import tkinter as tk
import webbrowser
from dataclasses import dataclass, field
from datetime import datetime, timedelta, timezone
from enum import Enum
from http.cookies import SimpleCookie
from logging.handlers import RotatingFileHandler
from tkinter import filedialog, messagebox, ttk
from typing import Dict, List, Optional, Set, Tuple
from urllib.parse import parse_qs, unquote, urlencode, urljoin, urlparse
from urllib.robotparser import RobotFileParser

import customtkinter as ctk
import requests
from bs4 import BeautifulSoup

# ─────────────────────────────────────────────────────────────
#  Frivilliga integrationer
# ─────────────────────────────────────────────────────────────
try:
    import trafilatura
    HAS_TRAFILATURA = True
except ImportError:
    HAS_TRAFILATURA = False

try:
    from playwright.async_api import async_playwright
    HAS_PLAYWRIGHT = True
except ImportError:
    HAS_PLAYWRIGHT = False

try:
    import aiohttp
    from aiohttp import CookieJar
    from yarl import URL as YarlURL
    HAS_AIOHTTP = True
except ImportError:
    HAS_AIOHTTP = False

try:
    import aiosqlite
    HAS_AIOSQLITE = True
except ImportError:
    HAS_AIOSQLITE = False

try:
    import uvloop
    HAS_UVLOOP = True
except ImportError:
    HAS_UVLOOP = False

try:
    import brotli  # noqa: F401  -- aiohttp upptäcker den automatiskt
    HAS_BROTLI = True
except ImportError:
    HAS_BROTLI = False

try:
    import fitz as pymupdf  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import load_workbook as xl_load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

ctk.set_appearance_mode("Light")
ctk.set_default_color_theme("blue")

# ─────────────────────────────────────────────────────────────
#  ENUMS OCH DATACLASSES
# ─────────────────────────────────────────────────────────────
class LogLevel(Enum):
    DEBUG = logging.DEBUG
    INFO = logging.INFO
    WARNING = logging.WARNING
    ERROR = logging.ERROR


class CrawlPriority(Enum):
    CRITICAL = 1
    SITEMAP = 5
    HIGH = 10
    MEDIUM = 15
    LOW = 20


class CrawlerState(Enum):
    IDLE = 0
    RUNNING = 1
    PAUSED = 2
    STOPPED = 3


@dataclass
class CrawlStats:
    pages_visited: int = 0
    pages_unchanged: int = 0
    pages_not_modified_304: int = 0   # Räknar 304-träffar separat
    pages_failed: int = 0
    playwright_fallbacks: int = 0
    documents_downloaded: int = 0
    bytes_downloaded: int = 0
    start_time: datetime = field(default_factory=datetime.now)
    end_time: Optional[datetime] = None

    @property
    def duration(self) -> timedelta:
        end = self.end_time or datetime.now()
        return end - self.start_time

    @property
    def pages_per_second(self) -> float:
        secs = self.duration.total_seconds()
        return self.pages_visited / secs if secs > 0 else 0.0


# ─────────────────────────────────────────────────────────────
#  HJÄLPFUNKTIONER & CHUNKING
# ─────────────────────────────────────────────────────────────
def slugify(text: str) -> str:
    text = str(text).replace('å', 'a').replace('ä', 'a').replace('ö', 'o')
    text = text.replace('Å', 'A').replace('Ä', 'A').replace('Ö', 'O')
    text = re.sub(r'[^\w\s-]', '', text)
    return re.sub(r'[-\s]+', '-', text).strip('-').lower()


# Standard-parametrar som aldrig påverkar sidans innehåll och därför ska strippas
# före URL-deduplicering. Lägg till nya här när du upptäcker dubletter i utmappen.
DEFAULT_IGNORE_QUERY_PARAMS = (
    # Generisk kampanj-/spårningsspårning
    'utm_source', 'utm_medium', 'utm_campaign', 'utm_term', 'utm_content',
    'fbclid', 'gclid', 'msclkid', 'mc_cid', 'mc_eid',
    'ref', 'source', 'igshid', 'mibextid',
    # Session-/auth-tokens som varierar mellan besök men inte ändrar innehåll
    'sessionid', 'jsessionid', 'phpsessid',
    # Sitevision (svenska kommuners CMS — Tyresö m.fl.).
    # `state` (addBookmark/removeBookmark) och `logout=true` är de vanligaste
    # som sett till fyra dubletter av samma sida.
    'state', 'logout', 'printerfriendly',
)

# Prefix-matchade parametrar — alla nycklar som börjar med dessa strippas.
# Detta täcker hela Sitevision-familjen (sv.url, sv.target, sv.viewportname,
# sv.scrollTo, sv.13.svid12_*, etc.) utan att vi behöver räkna upp varje variant.
DEFAULT_IGNORE_QUERY_PREFIXES = (
    'sv.',          # Sitevision — alla interna parametrar
    '_hsenc',       # HubSpot tracking
    '_hsmi',        # HubSpot tracking
)


def normalize_url(url: str,
                  ignore_query_params: Optional[List[str]] = None,
                  ignore_query_prefixes: Optional[List[str]] = None) -> str:
    """Normaliserar URL för stabil deduplicering.

    Tar bort:
      - tracking-/session-parametrar via exakt matchning (utm_*, fbclid, ...)
      - CMS-interna parametrar via prefix-matchning (sv.* för Sitevision, ...)
      - /index.html → /
      - trailing slash (för konsekvens)
      - URL-fragment (#...)
    """
    if ignore_query_params is None:
        ignore_query_params = DEFAULT_IGNORE_QUERY_PARAMS
    if ignore_query_prefixes is None:
        ignore_query_prefixes = DEFAULT_IGNORE_QUERY_PREFIXES

    # Sänk till lowercase en gång för effektivitet
    ignore_set = frozenset(p.lower() for p in ignore_query_params)
    ignore_prefixes = tuple(p.lower() for p in ignore_query_prefixes)

    try:
        parsed = urlparse(url.strip())
        scheme = parsed.scheme.lower() or 'http'
        netloc = parsed.netloc.lower()
        path = parsed.path or '/'
        path = re.sub(r'/index\.(html|htm|php)$', '/', path, flags=re.IGNORECASE)
        if path != '/' and path.endswith('/'):
            path = path.rstrip('/')
        query_params = parse_qs(parsed.query, keep_blank_values=True)
        filtered = {
            k: sorted(v) for k, v in query_params.items()
            if k.lower() not in ignore_set
            and not k.lower().startswith(ignore_prefixes)
        }
        query_string = urlencode(sorted(filtered.items()), doseq=True) if filtered else ""
        normalized = f"{scheme}://{netloc}{path}"
        if query_string:
            normalized += f"?{query_string}"
        return normalized.split('#')[0]
    except Exception:
        return url.strip()


def get_clean_hash(text: str) -> str:
    """Stabil hash av text efter normalisering av whitespace."""
    clean_text = re.sub(r'\s+', '', text).lower()
    return hashlib.sha256(clean_text.encode('utf-8')).hexdigest()


def stable_filename(url: str, save_format: str) -> str:
    """Filnamn som är stabilt mellan körningar (baserat på URL, inte content)."""
    parsed = urlparse(url)
    path_slug = slugify(parsed.path)[:60] or "root"
    url_digest = hashlib.md5(url.encode('utf-8')).hexdigest()[:8]
    return f"{path_slug}_{url_digest}{save_format}"


def semantic_chunk_text(sections: List[Dict], max_words: int = 400,
                        overlap_words: int = 50,
                        source_url: Optional[str] = None) -> List[Dict]:
    """Chunkar strukturerade sektioner till {heading, content, url}-objekt.

    `source_url` injiceras i varje chunk under nyckeln `url`. Det är harmlöst
    om RAG-systemet ignorerar fältet, men oerhört nyttigt när modellen får
    flera relaterade chunks i kontexten samtidigt — då kan källan plockas
    direkt från chunken istället för att gissas från rotnivå-metadata.
    """
    if not sections:
        return []

    chunks = []
    for sec in sections:
        heading = sec.get("heading", "Huvudinnehåll")
        text = sec.get("text", "").strip()
        if not text:
            continue

        words = text.split()
        if len(words) <= max_words:
            chunks.append({"heading": heading, "content": text})
        else:
            part_num = 1
            current_words: List[str] = []
            for w in words:
                current_words.append(w)
                if len(current_words) >= max_words:
                    chunk_heading = heading if part_num == 1 else f"{heading} (del {part_num})"
                    chunks.append({"heading": chunk_heading,
                                   "content": " ".join(current_words)})
                    part_num += 1
                    current_words = current_words[-overlap_words:]

            if len(current_words) > overlap_words:
                chunk_heading = heading if part_num == 1 else f"{heading} (del {part_num})"
                chunks.append({"heading": chunk_heading,
                               "content": " ".join(current_words)})

    total = len(chunks)
    for i, chunk in enumerate(chunks):
        chunk["chunk_index"] = i + 1
        chunk["total_chunks"] = total
        if source_url:
            chunk["url"] = source_url
    return chunks


def absolutize_markdown_links(text: str, base_url: str) -> str:
    """Gör relativa markdown-länkar [text](href) absoluta."""
    def _fix_md_link(m):
        href = m.group(2)
        if href and not href.startswith(('http://', 'https://', 'mailto:', '#')):
            href = urljoin(base_url, href)
        return f"[{m.group(1)}]({href})"
    return re.sub(r'\[([^\]]*)\]\(([^)]+)\)', _fix_md_link, text)


# ── Boilerplate/CMS-chrome som ska rensas från extraherad brödtext ──
# Sitevision (Tyresö m.fl.) injicerar feedback-widget, kontaktfooter och
# "Sidan publicerad av"-blocket i <main> — trafilatura/BS4 kan inte skilja
# detta från riktigt innehåll.
_CMS_BOILERPLATE_RE = re.compile(
    r'(?:'
    # Sitevision feedback-widget
    r'Tack för din medverkan!'
    r'|Du har hjälpt oss att förbättra webbplatsen'
    r'|Någonting gick fel\.?\s*Prova igen senare\.?'
    # Publiceringsinfo ("Sidan publicerad av:" + e-post/namn + datum)
    r'|\*{0,2}Sidan publicerad av:\*{0,2}.*'
    r'|\*{0,2}Senast uppdaterad:\*{0,2}.*'
    # Generisk "Var informationen till nytta?"-widget
    r'|Var informationen till nytta\??'
    r'|Skriv inte personuppgifter här'
    r'|Fältet är obligatoriskt'
    r')'
    r'\s*$',
    re.MULTILINE | re.IGNORECASE
)


def strip_cms_boilerplate(text: str) -> str:
    """Ta bort CMS-chrome/boilerplate (Sitevision m.fl.) från extraherad text.

    Rensningen sker radvis med regex. Rader som matchar kända boilerplate-
    mönster ersätts med tomrader, och sedan städas överflödiga tomrader.
    """
    cleaned = _CMS_BOILERPLATE_RE.sub('', text)
    # Kolla om rader med bara e-postlänk i slutet troligen är del av
    # "Sidan publicerad av"-blocket
    cleaned = re.sub(
        r'\n\[?[\w.+-]+@[\w.-]+\]?\(mailto:[^)]+\)\s*$',
        '', cleaned
    )
    # Städa upp: max två tomrader i rad
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    return cleaned.strip()


def downgrade_body_h1(text: str) -> str:
    """Ändra alla # (H1) till ## (H2) i brödtext.

    Crawlern sätter redan sin egen "# Titel" i filen, så H1-rubriker
    i den extraherade brödtexten blir dubbletter. Nedgradering till H2
    bevarar den visuella strukturen utan att förvirra RAG-chunkning.
    """
    return re.sub(r'^# ', '## ', text, flags=re.MULTILINE)


# ─────────────────────────────────────────────────────────────
#  DATABAS  (utökad med ETag / Last-Modified för conditional GET)
# ─────────────────────────────────────────────────────────────
class AsyncCrawlDatabase:
    """Cache med stöd för ETag och Last-Modified för 304-respons.

    Batched commits: ändringar samlas och commitas var COMMIT_BATCH_SIZE skrivning
    eller var COMMIT_BATCH_SECONDS sekund — vilket som kommer först. Detta
    eliminerar en fsync per sida.
    """

    COMMIT_BATCH_SIZE = 25
    COMMIT_BATCH_SECONDS = 5.0

    def __init__(self, db_path: str):
        self.db_path = db_path
        self.conn = None
        self._pending_writes = 0
        self._last_commit = time.monotonic()
        self._commit_lock = asyncio.Lock()

    async def connect(self):
        self.conn = await aiosqlite.connect(self.db_path)
        await self.conn.execute("PRAGMA journal_mode=WAL")
        await self.conn.execute("PRAGMA synchronous=NORMAL")  # snabbare, fortfarande crash-säkert i WAL
        await self._init_db()

    async def _init_db(self):
        await self.conn.execute('''
            CREATE TABLE IF NOT EXISTS page_cache (
                url TEXT PRIMARY KEY,
                content_hash TEXT,
                title TEXT,
                crawled_at TEXT,
                content_length INTEGER,
                etag TEXT,
                last_modified TEXT
            )
        ''')
        # Migration från äldre schema (lägger bara till om de saknas)
        async with self.conn.execute("PRAGMA table_info(page_cache)") as cur:
            cols = {row[1] for row in await cur.fetchall()}
        if 'etag' not in cols:
            await self.conn.execute("ALTER TABLE page_cache ADD COLUMN etag TEXT")
        if 'last_modified' not in cols:
            await self.conn.execute("ALTER TABLE page_cache ADD COLUMN last_modified TEXT")
        await self.conn.commit()

    async def get_cache(self, url: str) -> Optional[Dict]:
        async with self.conn.execute(
            "SELECT content_hash, title, crawled_at, content_length, etag, last_modified "
            "FROM page_cache WHERE url = ?", (url,)
        ) as cursor:
            row = await cursor.fetchone()
            if row:
                return {
                    'hash': row[0],
                    'title': row[1],
                    'crawled_at': row[2],
                    'content_length': row[3],
                    'etag': row[4],
                    'last_modified': row[5],
                }
        return None

    async def save_cache(self, url: str, content_hash: str, title: str,
                         length: int, etag: Optional[str] = None,
                         last_modified: Optional[str] = None):
        await self.conn.execute(
            'INSERT OR REPLACE INTO page_cache '
            '(url, content_hash, title, crawled_at, content_length, etag, last_modified) '
            'VALUES (?, ?, ?, ?, ?, ?, ?)',
            (url, content_hash, title, datetime.now().isoformat(),
             length, etag, last_modified)
        )
        await self._maybe_commit()

    async def touch_cache(self, url: str, etag: Optional[str] = None,
                          last_modified: Optional[str] = None):
        """Uppdatera bara crawled_at (och ev. ETag-värden) — för 304 Not Modified."""
        params = [datetime.now().isoformat()]
        sql = "UPDATE page_cache SET crawled_at=?"
        if etag is not None:
            sql += ", etag=?"
            params.append(etag)
        if last_modified is not None:
            sql += ", last_modified=?"
            params.append(last_modified)
        sql += " WHERE url=?"
        params.append(url)
        await self.conn.execute(sql, params)
        await self._maybe_commit()

    async def _maybe_commit(self):
        async with self._commit_lock:
            self._pending_writes += 1
            elapsed = time.monotonic() - self._last_commit
            if (self._pending_writes >= self.COMMIT_BATCH_SIZE
                    or elapsed >= self.COMMIT_BATCH_SECONDS):
                await self.conn.commit()
                self._pending_writes = 0
                self._last_commit = time.monotonic()

    async def flush(self):
        async with self._commit_lock:
            if self._pending_writes > 0:
                await self.conn.commit()
                self._pending_writes = 0
                self._last_commit = time.monotonic()

    async def get_all_records(self):
        async with self.conn.execute(
            "SELECT url, title, crawled_at, content_hash FROM page_cache "
            "ORDER BY crawled_at DESC"
        ) as cursor:
            return await cursor.fetchall()

    async def close(self):
        if self.conn:
            await self.flush()
            await self.conn.close()


# ─────────────────────────────────────────────────────────────
#  RATE LIMITER & QUEUE
# ─────────────────────────────────────────────────────────────
class PerDomainRateLimiter:
    """Async-vänlig per-domän rate limiter."""

    def __init__(self, requests_per_second: float):
        self.delay = 1.0 / requests_per_second if requests_per_second > 0 else 0
        self.last_requests: Dict[str, datetime] = {}
        self._lock = asyncio.Lock()

    async def async_wait(self, domain: str):
        async with self._lock:
            now = datetime.now()
            last_req = self.last_requests.get(domain, datetime.min)
            elapsed = (now - last_req).total_seconds()
            if elapsed < self.delay:
                sleep_time = self.delay - elapsed
                self.last_requests[domain] = now + timedelta(seconds=sleep_time)
            else:
                sleep_time = 0
                self.last_requests[domain] = now
        if sleep_time > 0:
            await asyncio.sleep(sleep_time)


class PriorityURLQueue:
    """Prioritetskö som rangordnar sidor efter "intressanthet"."""

    def __init__(self):
        self.queue: queue.PriorityQueue = queue.PriorityQueue()
        self.seen_urls: Set[str] = set()
        self._lock = threading.Lock()

        self.boost_words = ["policy", "om-oss", "kontakt", "regler", "guide"]
        self.penalty_words = ["nyheter", "arkiv", "blogg", "kalender", "202"]

    def add_url(self, url: str, depth: int = 0,
                base_priority: int = CrawlPriority.MEDIUM.value) -> bool:
        normalized = normalize_url(url)
        with self._lock:
            if normalized in self.seen_urls:
                return False

            score = base_priority
            lower_url = normalized.lower()
            if any(w in lower_url for w in self.boost_words):
                score -= 3
            if any(w in lower_url for w in self.penalty_words):
                score += 5
            score += depth

            self.queue.put((score, time.time(), depth, normalized))
            self.seen_urls.add(normalized)
            return True

    def get_next(self) -> Optional[Tuple[int, str]]:
        try:
            item = self.queue.get_nowait()
            return (item[2], item[3])
        except queue.Empty:
            return None

    def size(self) -> int:
        return self.queue.qsize()


# ─────────────────────────────────────────────────────────────
#  LOGIN-DETEKTOR
# ─────────────────────────────────────────────────────────────
class LoginDetector:
    """Avgör om en respons faktiskt representerar en utgången session.

    Tre lager av signaler, från starkast till svagast:
      1. URL-redirect: slut-URL:en innehåller en känd login-sökväg
         (t.ex. /login, /cas/login, /adfs/ls, ?SAMLRequest=, /saml2/sso)
      2. HTTP-status: 401 / 403
      3. Innehåll: små HTML-sidor som ser ut som login-formulär
         (få forms, exakt ett password-fält, eller känd login-text)

    SAML-matchning görs på URL-mönster (/saml/, /saml2/) istället för
    substrängar i brödtext — "saml" förekommer i svenska ord som *samla,
    samling, samlade* och ger annars falska positiva. Password-fält kräver
    att sidan är liten (<15 kB) med max 2 forms för att undvika att sidor
    med login-widget i headern flaggas felaktigt.
    """

    DEFAULT_URL_PATTERNS = (
        '/login', '/signin', '/sign-in', '/log-in',
        '/cas/login', '/adfs/ls', '/oauth/authorize',
        '/saml/', '/saml2/', '/sso/',
        'samlrequest=', 'returnurl=', 'redirect_uri=',
        'logon.aspx', '/auth/realms/',
    )

    STRONG_CONTENT_SIGNALS = (
        'id="loginform"', "id='loginform'",
        'class="login-form"', 'class="loginform"',
        'inloggning krävs', 'du måste logga in',
        'din session har gått ut', 'session expired',
        'please sign in to continue', 'sign in to continue',
        'authentication required',
    )

    def __init__(self, extra_url_patterns: Optional[List[str]] = None,
                 max_login_html_size: int = 30000):
        self.url_patterns = list(self.DEFAULT_URL_PATTERNS)
        if extra_url_patterns:
            self.url_patterns.extend(p.lower() for p in extra_url_patterns)
        self.max_login_html_size = max_login_html_size

    def is_login_redirect(self, original_url: str, final_url: Optional[str]) -> bool:
        """Stark signal: sidan redirecterade till en login-URL."""
        if not final_url:
            return False
        # Bara intressant om vi faktiskt redirectades NÅGONSTANS, eller om slut-URL:en
        # tydligt är en login-sida även utan redirect.
        final_lower = final_url.lower()
        return any(p in final_lower for p in self.url_patterns)

    def is_login_status(self, http_status: Optional[int]) -> bool:
        return http_status in (401, 407)

    def is_login_content(self, html: Optional[str]) -> bool:
        """Innehållsbaserad heuristik. Endast små HTML-sidor undersöks
        för att undvika falska positiva på vanliga innehållssidor."""
        if not html:
            return False
        if len(html) > self.max_login_html_size:
            return False  # Innehållssidor är typiskt 50+ kB; login-sidor är små
        html_lower = html.lower()

        # Stark signal: explicit login-text eller login-form-id
        if any(s in html_lower for s in self.STRONG_CONTENT_SIGNALS):
            return True

        # Svagare: ett ENDA password-fält i en liten sida med få forms.
        # Vi kräver att det är en kort sida (<15 kB) och har max 2 forms.
        if 'type="password"' in html_lower or "type='password'" in html_lower:
            if len(html) < 15000:
                form_count = html_lower.count('<form')
                pw_count = (html_lower.count('type="password"')
                            + html_lower.count("type='password'"))
                if form_count <= 2 and pw_count == 1:
                    return True
        return False

    def detect(self, original_url: str, final_url: Optional[str],
               http_status: Optional[int], html: Optional[str]) -> bool:
        if self.is_login_redirect(original_url, final_url):
            return True
        if self.is_login_status(http_status):
            return True
        if self.is_login_content(html):
            return True
        return False


# ─────────────────────────────────────────────────────────────
#  DOKUMENT-MANIFEST
# ─────────────────────────────────────────────────────────────
class DocumentManifest:
    """Spårar nedladdade dokument och var de länkades från.

    Producerar en `manifest.json` i utmappen som låter ett RAG-system koppla
    en PDF-text tillbaka till sin ursprungssida på intranätet. Utan detta
    blir varje PDF en "isolerad ö" — modellen kan citera innehåll men inte
    säga vilken intranätsida som beskriver dokumentet.

    Strukturen för varje dokument:
      filename       — det faktiska filnamnet på disk
      download_url   — direktlänken som crawlern hämtade
      referer_url    — sidan på intranätet som länkade till dokumentet
      referer_title  — titel på den länkande sidan
      link_text      — den klickbara textens innehåll (ofta beskrivande)
      size_bytes     — filstorlek
      downloaded_at  — ISO-tidsstämpel
      additional_referers — om PDFen länkas från flera sidor
    """

    def __init__(self):
        # download_url → list[{referer_url, referer_title, link_text, found_at}]
        self._referers: Dict[str, List[Dict[str, str]]] = {}
        # download_url → {filename, size_bytes, downloaded_at}
        self._downloads: Dict[str, Dict] = {}
        self._lock = asyncio.Lock()

    async def record_link(self, doc_url: str, referer_url: str,
                          referer_title: str, link_text: str):
        """Anropas i process_page när en länk till ett dokument hittas."""
        entry = {
            "referer_url": referer_url,
            "referer_title": referer_title or "",
            "link_text": (link_text or "").strip(),
            "found_at": datetime.now().isoformat(),
        }
        async with self._lock:
            referer_list = self._referers.setdefault(doc_url, [])
            # Dedupa: samma referer_url ska inte registreras två gånger
            if not any(r["referer_url"] == referer_url for r in referer_list):
                referer_list.append(entry)

    async def record_download(self, doc_url: str, filename: str,
                              size_bytes: int):
        """Anropas i download_document efter framgångsrik nedladdning."""
        async with self._lock:
            self._downloads[doc_url] = {
                "filename": filename,
                "size_bytes": size_bytes,
                "downloaded_at": datetime.now().isoformat(),
            }

    def build(self, domain: str) -> Dict:
        """Producerar slutlig manifest-struktur. Kallas en gång vid crawl-slut."""
        documents = []
        # Sortera på filnamn för stabil output mellan körningar
        for doc_url in sorted(self._downloads.keys(),
                              key=lambda u: self._downloads[u]["filename"]):
            dl_info = self._downloads[doc_url]
            referers = self._referers.get(doc_url, [])

            entry = {
                "filename": dl_info["filename"],
                "download_url": doc_url,
                "size_bytes": dl_info["size_bytes"],
                "downloaded_at": dl_info["downloaded_at"],
            }

            if referers:
                primary = referers[0]
                entry["referer_url"] = primary["referer_url"]
                entry["referer_title"] = primary["referer_title"]
                entry["link_text"] = primary["link_text"]
                if len(referers) > 1:
                    entry["additional_referers"] = referers[1:]
            documents.append(entry)

        # Inkludera även dokument som länkades men inte laddades ner
        # (kan hända vid avbruten crawl) — bra för felsökning
        orphan_links = []
        for doc_url, refs in self._referers.items():
            if doc_url not in self._downloads:
                orphan_links.append({
                    "download_url": doc_url,
                    "referers": refs,
                })

        manifest = {
            "generated_at": datetime.now().isoformat(),
            "domain": domain,
            "document_count": len(documents),
            "documents": documents,
        }
        if orphan_links:
            manifest["orphan_links"] = orphan_links
            manifest["orphan_count"] = len(orphan_links)
        return manifest


# ─────────────────────────────────────────────────────────────
#  DOKUMENT → MARKDOWN KONVERTERARE
# ─────────────────────────────────────────────────────────────
class DocumentConverter:
    """Konverterar binära dokument (PDF, Word, Excel, PPTX) till Markdown.

    Extraherar text och skapar en .md-fil med intranät-URL:en i toppen,
    så att RAG-system kan citera rätt källa även för dokument-text.
    Originaldokumentet behålls alltid i dokument/-mappen.
    """

    SUPPORTED = {
        '.pdf': 'PDF', '.docx': 'Word', '.dotx': 'Word-mall',
        '.xlsx': 'Excel', '.pptx': 'PowerPoint',
    }
    # Äldre format (.doc, .xls) stöds ej av Python-biblioteken
    UNSUPPORTED_LEGACY = {'.doc', '.xls'}

    # Generiska länktexter som inte är beskrivande nog att använda som
    # dokument-titel. Matchning sker case-insensitivt efter strip().
    GENERIC_LINK_TEXTS = {
        'ladda ner', 'ladda ner fil', 'ladda ned', 'ladda ned fil',
        'hämta', 'hämta fil', 'download', 'download file',
        'öppna', 'öppna fil', 'öppna dokument',
        'visa', 'visa fil', 'visa dokument',
        'klicka här', 'läs mer', 'read more', 'click here',
        'länk', 'link', 'pdf', 'dokument', 'document',
        '(via sitemap)', 'via sitemap',
    }

    def __init__(self, output_dir: str,
                 log_fn=None, pii_cleaner=None):
        self.texts_dir = os.path.join(output_dir, "texter")
        self._log = log_fn
        self._clean_pii = pii_cleaner
        os.makedirs(self.texts_dir, exist_ok=True)

    def can_convert(self, filepath: str) -> bool:
        ext = os.path.splitext(filepath)[1].lower()
        if ext in self.UNSUPPORTED_LEGACY:
            return False
        if ext == '.pdf' and not HAS_PYMUPDF:
            return False
        if ext in ('.docx', '.dotx') and not HAS_DOCX:
            return False
        if ext == '.xlsx' and not HAS_OPENPYXL:
            return False
        if ext == '.pptx' and not HAS_PPTX:
            return False
        return ext in self.SUPPORTED

    def convert(self, filepath: str, doc_url: str,
                referer_url: str = "", referer_title: str = "",
                link_text: str = "") -> Optional[str]:
        """Konverterar dokument till .md. Returnerar sökväg eller None.

        Metadata bäddas in i brödtexten (inte bara i header) eftersom RAG-
        pipelines ofta kapar de första raderna före retrieval. Källan
        upprepas också sist i filen — om Sveas chunking splittrar dokumentet
        i flera bitar säkrar det att åtminstone en chunk har källan med.
        """
        ext = os.path.splitext(filepath)[1].lower()
        extractors = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.dotx': self._extract_docx,
            '.xlsx': self._extract_xlsx,
            '.pptx': self._extract_pptx,
        }
        extractor = extractors.get(ext)
        if not extractor:
            return None

        try:
            text = extractor(filepath)
            if not text or len(text.strip()) < 20:
                if self._log:
                    self._log(
                        f"  ⚠ Konvertering gav för lite text "
                        f"({len(text.strip()) if text else 0} tecken): "
                        f"{os.path.basename(filepath)}", LogLevel.WARNING)
                return None

            if self._clean_pii:
                text = self._clean_pii(text)

            # Titel-prioritet (faller från bäst till sista utvägen):
            #   1. link_text   — <a>-taggens text (t.ex. "Eko ack 2025 Januari"),
            #                    kräver >5 tecken OCH att texten inte är generisk
            #                    ("Ladda ner fil", "Download", etc.)
            #   2. metadata_title — dokumentets egen title (PDF/Office metadata)
            #   3. humaniserat filnamn ("Eko ack 2025 januari")
            #   4. referer_title — sidans <title> (ofta för generisk,
            #                      t.ex. "Startsida - Tyresö kommun")
            #   5. rå slug som sista utväg
            filename = os.path.basename(filepath)
            base_name = os.path.splitext(filename)[0]
            metadata_title = self._extract_doc_metadata_title(filepath, ext)
            humanized = self._humanize_filename(base_name)
            clean_link = (link_text or "").strip()

            # De-duplicera dubbla länktexter ("Ladda ner fil Ladda ner fil")
            if clean_link:
                total_len = len(clean_link)
                if total_len % 2 == 1:  # udda längd → testa med mellanslag i mitten
                    mid = total_len // 2
                    if (clean_link[mid] == ' '
                            and clean_link[:mid] == clean_link[mid + 1:]):
                        clean_link = clean_link[:mid]
                elif total_len % 2 == 0:
                    mid = total_len // 2
                    if clean_link[:mid].rstrip() == clean_link[mid:].lstrip():
                        clean_link = clean_link[:mid].rstrip()

            def _is_usable_link_text(text: str) -> bool:
                """True om link_text är tillräckligt beskrivande för titel."""
                if len(text) <= 5:
                    return False
                return text.lower() not in self.GENERIC_LINK_TEXTS

            raw_title = (
                (clean_link if _is_usable_link_text(clean_link) else "")
                or metadata_title
                or humanized
                or (referer_title.strip() if referer_title else "")
                or base_name
            )

            # Strippa storleks-suffix som "pdf, 535.8 kB, öppnas i nytt fönster"
            title = re.sub(
                r'\s*(pdf|docx?|xlsx?|pptx?|pages?)?\s*,?\s*'
                r'\d+(?:[.,]\d+)?\s*[kKmMgG]?[bB][, .].*$',
                '', raw_title, flags=re.IGNORECASE).strip()
            title = re.sub(r'\s*[, .]+\s*$', '', title) or raw_title

            file_type = self.SUPPORTED.get(ext, "Dokument")

            # Bygg metadata-block som SYNLIG brödtext (fetstil), inte ren header.
            # Sveas chunker skippar typiskt header-text men behåller brödtext.
            meta_lines = [f"# {title}", ""]
            if referer_url:
                meta_lines.append(f"**Källa:** {referer_url}")
            if doc_url:
                meta_lines.append(f"**Dokument-URL:** {doc_url}")
            if filename:
                # Visa originalfilnamnet (utan vår interna hash-suffix)
                display_name = re.sub(r'_[0-9a-f]{6}(\.\w+)$', r'\1', filename)
                meta_lines.append(f"**Filnamn:** {display_name}")
            meta_lines.append(f"**Filtyp:** {file_type}")
            meta_lines.append("")
            meta_lines.append("---")
            meta_lines.append("")

            # Upprepa källan sist i filen — om dokumentet chunkas i flera
            # bitar har åtminstone första och sista bitarna källinformation.
            footer_lines = ["", "", "---", ""]
            if referer_url:
                footer_lines.append(f"**Källa:** {referer_url}")
            if doc_url:
                footer_lines.append(f"**Dokument-URL:** {doc_url}")

            text = downgrade_body_h1(text)
            md_content = ("\n".join(meta_lines) + text
                          + "\n".join(footer_lines))

            # Filnamn: Återanvänd den hash som download_document redan satt
            # på safe_filename (formatet "<slug>_<hash>.<ext>"). Om hashen
            # redan finns i filnamnet undviks dubbla hash-suffix.
            existing_hash_match = re.search(r'_([0-9a-f]{6})$', base_name)
            if existing_hash_match:
                # Hashen finns redan i filnamnet — återanvänd den
                base = base_name[:50]  # behåll _ddf505 som det är
                md_filename = f"{base}_doc.md"
            else:
                # Filnamnet saknar hash (skickat in från annan väg) — generera en
                base = slugify(base_name)[:50]
                url_hash = hashlib.md5(doc_url.encode('utf-8')).hexdigest()[:6]
                md_filename = f"{base}_{url_hash}_doc.md"
            md_path = os.path.join(self.texts_dir, md_filename)

            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            return md_path

        except Exception as e:
            if self._log:
                self._log(
                    f"  ✗ Konvertering misslyckades "
                    f"({os.path.basename(filepath)}): {e}", LogLevel.ERROR)
            return None

    # ─── Titelextraktion ────────────────────────────────────
    @staticmethod
    def _humanize_filename(name: str) -> str:
        """Gör ett slug-namn läsbart: 'ansokan-om-x_ad7919' → 'Ansokan om x'.

        Eftersom slugify har strippat svenska tecken (ä→a, ö→o) går det
        inte att återskapa originalnamnet helt — men "Ansokan om
        parkeringstillstand" är vida bättre än "ansokan-om-parkeringstillstand_ad7919".
        """
        # Ta bort hash-suffix
        name = re.sub(r'_[0-9a-f]{6}$', '', name)
        # Bindestreck och underscore → mellanslag
        name = re.sub(r'[-_]+', ' ', name).strip()
        # Stor första bokstav (ev. övriga ord lämnas som de är —
        # vi vet inte vilka som är egennamn)
        if name:
            name = name[0].upper() + name[1:]
        return name

    def _extract_doc_metadata_title(self, filepath: str, ext: str) -> str:
        """Plocka ut dokumenttitel ur filens egen metadata om sådan finns.

        Många PDF/Office-dokument har ett `title`-fält i sin metadata som
        ofta är mer beskrivande än filnamnet ("Ansökan om parkeringstillstånd
        för rörelsehindrad" vs "ansokan-om-parkeringstillstand").
        """
        try:
            if ext == '.pdf' and HAS_PYMUPDF:
                doc = pymupdf.open(filepath)
                title = (doc.metadata or {}).get('title', '')
                doc.close()
                return title.strip() if title else ""
            if ext in ('.docx', '.dotx') and HAS_DOCX:
                doc = DocxDocument(filepath)
                title = (doc.core_properties.title or '').strip()
                return title
            if ext == '.xlsx' and HAS_OPENPYXL:
                wb = xl_load_workbook(filepath, read_only=True, data_only=True)
                title = (wb.properties.title or '').strip()
                wb.close()
                return title
            if ext == '.pptx' and HAS_PPTX:
                prs = PptxPresentation(filepath)
                title = (prs.core_properties.title or '').strip()
                return title
        except Exception:
            pass
        return ""

    def _extract_pdf(self, filepath: str) -> str:
        doc = pymupdf.open(filepath)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                pages.append(f"## Sida {i + 1}\n\n{text.strip()}")
        doc.close()
        return "\n\n".join(pages)

    def _extract_docx(self, filepath: str) -> str:
        doc = DocxDocument(filepath)
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style and para.style.name and para.style.name.startswith('Heading'):
                try:
                    level = int(para.style.name.replace('Heading ', '')
                                .replace('Heading', '1'))
                except ValueError:
                    level = 2
                parts.append(f"{'#' * min(level + 1, 6)} {text}")
            else:
                parts.append(text)
        for table in doc.tables:
            rows = []
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip().replace('\n', ' ')
                         for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    rows.append("|" + "|".join(["---"] * len(cells)) + "|")
            if rows:
                parts.append("\n".join(rows))
        return "\n\n".join(parts)

    def _extract_xlsx(self, filepath: str) -> str:
        wb = xl_load_workbook(filepath, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## {sheet_name}")
            rows_data = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).replace('\n', ' ') if c is not None
                         else "" for c in row]
                if any(c for c in cells):
                    rows_data.append("| " + " | ".join(cells) + " |")
            if rows_data:
                col_count = rows_data[0].count("|") - 1
                rows_data.insert(1, "|" + "|".join(["---"] * col_count) + "|")
                parts.append("\n".join(rows_data))
        wb.close()
        return "\n\n".join(parts)

    def _extract_pptx(self, filepath: str) -> str:
        prs = PptxPresentation(filepath)
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(
                    f"## Bild {i + 1}\n\n" + "\n\n".join(slide_texts))
        return "\n\n".join(parts)


# ─────────────────────────────────────────────────────────────
#  ASYNC WEBB CRAWLER CORE
# ─────────────────────────────────────────────────────────────
# Returneras av nätverkslagret för konsekvent hantering uppströms
@dataclass
class FetchResult:
    body: Optional[bytes] = None         # Råa bytes (eller None om body inte hämtades)
    text: Optional[str] = None           # Avkodad text (om relevant content-type)
    content_type: str = ""
    final_url: str = ""                  # URL efter ev. redirects
    status: int = 0
    etag: Optional[str] = None
    last_modified: Optional[str] = None
    not_modified: bool = False           # True om servern svarade 304


# Filextensioner som indikerar binära dokument — för dessa gör vi HEAD först
DOCUMENT_EXTENSIONS = {
    '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx',
    '.odt', '.ods', '.odp', '.rtf', '.csv', '.zip',
}
# Filextensioner vi aldrig hämtar (bilder, fonts, video, etc.)
# .zip ligger MEDVETET inte här — den finns i DOCUMENT_EXTENSIONS, så
# zip-arkiv kan laddas ner som dokument om download_docs är på.
SKIP_EXTENSIONS = {
    '.jpg', '.jpeg', '.png', '.gif', '.svg', '.webp', '.ico', '.bmp', '.tiff',
    '.rar', '.exe', '.mp4', '.mp3', '.avi', '.mov', '.css', '.js',
    '.woff', '.woff2', '.ttf', '.otf', '.eot',
}


class AsyncWebCrawler:
    def __init__(self, config: dict, msg_queue: Optional[queue.Queue] = None):
        self.config = config
        self.msg_queue = msg_queue
        self.state = CrawlerState.RUNNING
        self.stats = CrawlStats()
        self.active_tasks = 0

        self.start_url = normalize_url(config["start_url"])
        self.output_dir = config["output_dir"]
        self.delay = config["delay"]
        self.max_pages = config["max_pages"]
        self.max_depth = config.get("max_depth", 0)
        self.save_format = config.get("save_format", ".md")
        self.use_hybrid = config.get("use_hybrid", True)
        self.use_trafilatura = config.get("use_trafilatura", HAS_TRAFILATURA)

        # Samtidighet: justerbar via config
        self.concurrency = max(1, int(config.get("concurrency", 10)))
        self.playwright_concurrency = max(1, int(config.get("playwright_concurrency", 2)))

        self.find_sitemap = config.get("find_sitemap", True)
        self.robot_parser: Optional[RobotFileParser] = None

        parsed_start = urlparse(self.start_url)
        self.domain = parsed_start.netloc.lower()
        self.base_url = f"{parsed_start.scheme}://{parsed_start.netloc}"

        os.makedirs(self.output_dir, exist_ok=True)
        self.db = AsyncCrawlDatabase(
            os.path.join(self.output_dir, f"{slugify(self.domain)}_cache.db")
        )
        self.url_queue = PriorityURLQueue()
        self.url_queue.add_url(self.start_url, depth=0,
                               base_priority=CrawlPriority.CRITICAL.value)

        self.rate_limiter = PerDomainRateLimiter(
            requests_per_second=1.0 / max(self.delay, 0.05)
        )
        self.downloaded_files: Set[str] = set()
        self.manifest = DocumentManifest()
        self.convert_docs = config.get("convert_docs_to_md", False)
        self.converter: Optional[DocumentConverter] = None
        if self.convert_docs:
            self.converter = DocumentConverter(
                self.output_dir,
                log_fn=self._log,
                pii_cleaner=self.clean_pii,
            )
        self.visited_sitemaps: Set[str] = set()
        self.login_event = threading.Event()
        self.saved_cookies: List[Dict] = []

        self.async_download_lock = asyncio.Lock()
        self.async_pw_lock = asyncio.Lock()
        self.async_stats_lock = asyncio.Lock()
        self.async_sitemap_lock = asyncio.Lock()

        self.semaphore = asyncio.Semaphore(self.concurrency)
        self.playwright_semaphore = asyncio.Semaphore(self.playwright_concurrency)

        # Login-detektor — bara aktiv när vi använder login-läge
        self.login_detector = LoginDetector()
        self._login_detection_enabled = (
            config.get("headless_mode") == "login_then_headless"
        )

        self.req_session: Optional[aiohttp.ClientSession] = None
        self._pw = None
        self._browser = None
        self._context = None

        self.crawl_session_id = datetime.now().strftime("%Y%m%d_%H%M%S")
        self.logger = logging.getLogger(f'Crawler_{id(self)}')
        self.logger.setLevel(logging.DEBUG)
        log_dir = os.path.join(self.output_dir, 'logs')
        os.makedirs(log_dir, exist_ok=True)
        fh = RotatingFileHandler(
            os.path.join(log_dir, f'crawl_{self.crawl_session_id}.log'),
            maxBytes=5 * 1024 * 1024, backupCount=2, encoding='utf-8'
        )
        fh.setFormatter(logging.Formatter(
            '%(asctime)s - %(levelname)s - %(message)s', datefmt='%H:%M:%S'
        ))
        self.logger.addHandler(fh)

        self._log(f"🚀 Initierar ASYNC crawl för {self.domain} (v7.0)")
        if self.convert_docs:
            libs = []
            if HAS_PYMUPDF: libs.append('PDF')
            if HAS_DOCX: libs.append('Word')
            if HAS_OPENPYXL: libs.append('Excel')
            if HAS_PPTX: libs.append('PPTX')
            self._log(f"✓ Dokument→Markdown aktiv ({', '.join(libs) or 'inga bibliotek!'})")
        if HAS_BROTLI:
            self._log("✓ Brotli-stöd aktivt", LogLevel.DEBUG)

    # ─── Logging & GUI ──────────────────────────────────────
    def _log(self, msg: str, level=LogLevel.INFO):
        if level == LogLevel.DEBUG:
            self.logger.debug(msg)
        elif level == LogLevel.INFO:
            self.logger.info(msg)
        elif level == LogLevel.WARNING:
            self.logger.warning(msg)
        elif level == LogLevel.ERROR:
            self.logger.error(msg)
        if self.msg_queue:
            self.msg_queue.put(("log", f"[{level.name}] {msg}"))
        else:
            print(f"[{level.name}] {msg}")

    def _gui_update(self, url: str, status: str, title: str):
        nya_eller_sparade = self.stats.pages_visited - self.stats.pages_unchanged
        queue_size = self.url_queue.size() + self.active_tasks
        pages_done = self.stats.pages_visited

        eta_str = "Beräknar..."
        if pages_done > 2:
            avg_time = self.stats.duration.total_seconds() / pages_done
            remaining = queue_size
            if self.max_pages > 0:
                remaining = min(queue_size, self.max_pages - pages_done)
            eta_sec = avg_time * remaining
            if remaining <= 0:
                eta_str = "Klar snart"
            elif eta_sec > 3600:
                eta_str = f"{int(eta_sec // 3600)}h {int((eta_sec % 3600) // 60)}m"
            elif eta_sec > 60:
                eta_str = f"{int(eta_sec // 60)}m {int(eta_sec % 60)}s"
            else:
                eta_str = f"{int(eta_sec)}s"

        if self.msg_queue:
            safe_title = title.replace('\x00', '') if title else "Ingen titel"
            self.msg_queue.put(("table", (url, status, safe_title)))
            self.msg_queue.put(("stats_data", (
                self.stats.pages_visited,
                nya_eller_sparade,
                self.stats.documents_downloaded,
                self.stats.pages_unchanged + self.stats.pages_not_modified_304,
                queue_size,
                self.stats.pages_failed,
                eta_str,
            )))

    # ─── URL-validering ─────────────────────────────────────
    def is_valid_url(self, url: str) -> bool:
        if len(url) > 2000:
            return False
        try:
            parsed = urlparse(url)
            if parsed.scheme not in ('http', 'https'):
                return False

            parsed_path = parsed.path.lower()
            ext = posixpath.splitext(parsed_path)[1]
            if ext in SKIP_EXTENSIONS:
                return False

            if any(p in parsed_path for p in ('/images/', '/media/', '/assets/')):
                if any(img in parsed_path
                       for img in ('.jpg', '.jpeg', '.png', '.gif', '.webp')):
                    return False

            domain_core = self.domain.replace('www.', '')
            link_domain = parsed.netloc.lower().replace('www.', '')
            if self.config.get("strict_domain", True) and link_domain != domain_core:
                return False

            lower_url = url.lower()
            for kw in self.config.get("exclude_keywords", []):
                if kw and kw in lower_url:
                    return False
            req_kws = self.config.get("require_keywords", [])
            if req_kws and not any(kw in lower_url for kw in req_kws):
                return False

            if self.robot_parser and not self.robot_parser.can_fetch('*', url):
                return False
            return True
        except Exception:
            return False

    # ─── Nätverkslager ──────────────────────────────────────
    async def _create_session(self) -> aiohttp.ClientSession:
        """Skapar aiohttp-session med riktig CookieJar och Brotli om tillgängligt.

        Cookies från Playwright-inloggningen filtreras till samma toppdomän
        innan de injiceras i CookieJar — undviker att tappa secure/httponly-flaggor.
        """
        connector = aiohttp.TCPConnector(
            limit=max(20, self.concurrency * 2),
            limit_per_host=self.concurrency,
            ttl_dns_cache=300,
        )
        accept_encoding = "gzip, deflate"
        if HAS_BROTLI:
            accept_encoding = "br, " + accept_encoding

        jar = CookieJar(unsafe=True)  # tillåt även IP-baserade cookies
        headers = {
            'User-Agent': ('Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
                           'AppleWebKit/537.36 (KHTML, like Gecko) '
                           'Chrome/124.0.0.0 Safari/537.36'),
            'Accept-Encoding': accept_encoding,
            'Accept': ('text/html,application/xhtml+xml,application/xml;q=0.9,'
                       'image/avif,image/webp,*/*;q=0.8'),
            'Accept-Language': 'sv-SE,sv;q=0.9,en;q=0.8',
        }

        session = aiohttp.ClientSession(
            headers=headers,
            connector=connector,
            cookie_jar=jar,
            timeout=aiohttp.ClientTimeout(total=30, connect=10),
        )

        # Injicera cookies från Playwright (om vi har sådana från login-läget)
        if self.saved_cookies:
            self._inject_playwright_cookies(jar, self.saved_cookies)
        return session

    def _inject_playwright_cookies(self, jar: CookieJar, cookies: List[Dict]):
        """Konverterar Playwright-cookies till format som aiohttp's CookieJar förstår."""
        added = 0
        for c in cookies:
            try:
                domain = (c.get('domain') or '').lstrip('.')
                if not domain:
                    domain = self.domain
                # Bygg en URL för update_cookies så att domain/path bevaras
                scheme = 'https' if c.get('secure') else 'http'
                cookie_url = f"{scheme}://{domain}{c.get('path', '/')}"

                sc = SimpleCookie()
                sc[c['name']] = c['value']
                m = sc[c['name']]
                if c.get('path'):
                    m['path'] = c['path']
                if c.get('expires') and c['expires'] > 0:
                    # Konvertera epoch till GMT-sträng
                    try:
                        expires_dt = datetime.fromtimestamp(c['expires'], tz=timezone.utc)
                        m['expires'] = expires_dt.strftime("%a, %d %b %Y %H:%M:%S GMT")
                    except Exception:
                        pass
                if c.get('secure'):
                    m['secure'] = True
                if c.get('httpOnly'):
                    m['httponly'] = True

                jar.update_cookies(sc, response_url=YarlURL(cookie_url))
                added += 1
            except Exception as e:
                self._log(f"Kunde inte överföra cookie {c.get('name')}: {e}",
                          LogLevel.DEBUG)
        self._log(f"✓ Överförde {added} cookies från login-sessionen")

    async def fetch(self, url: str, method: str = 'GET',
                    cached: Optional[Dict] = None,
                    max_retries: int = 3,
                    decode_text: bool = True) -> Optional[FetchResult]:
        """Enhetlig nätverkshämtning med retries, conditional GET och slut-URL.

        - method='HEAD' → bara content-type
        - cached + ETag/Last-Modified → conditional GET → 304 returneras som
          FetchResult(not_modified=True)
        - decode_text=True → returnerar str i .text, annars bara bytes i .body
        """
        base_delay = 1.0

        # Bygg conditional-headers från cache
        extra_headers: Dict[str, str] = {}
        if method == 'GET' and cached:
            if cached.get('etag'):
                extra_headers['If-None-Match'] = cached['etag']
            if cached.get('last_modified'):
                extra_headers['If-Modified-Since'] = cached['last_modified']

        for attempt in range(max_retries + 1):
            # Respektera pause/stop mellan försök
            while self.state == CrawlerState.PAUSED:
                await asyncio.sleep(0.3)
            if self.state == CrawlerState.STOPPED:
                return None

            try:
                if method == 'HEAD':
                    async with self.req_session.head(
                        url, allow_redirects=True,
                        timeout=aiohttp.ClientTimeout(total=10)
                    ) as resp:
                        if resp.status in (403, 405):
                            # Servern stödjer inte HEAD — låt anroparen falla tillbaka till GET
                            return FetchResult(
                                content_type=resp.headers.get('Content-Type', '').lower(),
                                final_url=str(resp.url), status=resp.status,
                            )
                        if resp.status in (429, 500, 502, 503, 504):
                            if attempt == max_retries:
                                return None
                            await asyncio.sleep(base_delay * (2 ** attempt))
                            continue
                        if resp.status >= 400:
                            return None
                        return FetchResult(
                            content_type=resp.headers.get('Content-Type', '').lower(),
                            final_url=str(resp.url),
                            status=resp.status,
                        )

                # GET
                async with self.req_session.get(
                    url, headers=extra_headers,
                    timeout=aiohttp.ClientTimeout(total=20)
                ) as resp:
                    if resp.status == 304:
                        return FetchResult(
                            final_url=str(resp.url), status=304,
                            not_modified=True,
                            etag=resp.headers.get('ETag'),
                            last_modified=resp.headers.get('Last-Modified'),
                        )
                    if resp.status in (429, 500, 502, 503, 504):
                        if attempt == max_retries:
                            return None
                        await asyncio.sleep(base_delay * (2 ** attempt))
                        continue
                    if resp.status >= 400 and resp.status not in (401, 403):
                        # 401/403 returneras till anroparen så login-detektorn kan agera
                        return None

                    content_type = resp.headers.get('Content-Type', '').lower()
                    final_url = str(resp.url)
                    etag = resp.headers.get('ETag')
                    last_mod = resp.headers.get('Last-Modified')

                    # Avgör om vi ska läsa body
                    is_textual = (
                        decode_text and (
                            'text' in content_type
                            or 'html' in content_type
                            or 'json' in content_type
                            or 'xml' in content_type
                            or content_type == ''
                        )
                    )
                    if is_textual:
                        try:
                            text = await resp.text(errors='replace')
                        except UnicodeDecodeError:
                            raw = await resp.read()
                            text = raw.decode('utf-8', errors='replace')
                        return FetchResult(
                            text=text, content_type=content_type,
                            final_url=final_url, status=resp.status,
                            etag=etag, last_modified=last_mod,
                        )
                    else:
                        body = await resp.read()
                        return FetchResult(
                            body=body, content_type=content_type,
                            final_url=final_url, status=resp.status,
                            etag=etag, last_modified=last_mod,
                        )
            except asyncio.CancelledError:
                raise
            except (aiohttp.ClientError, asyncio.TimeoutError) as e:
                if attempt == max_retries:
                    self._log(f"  ✗ Nätverksfel: {url} ({e})", LogLevel.DEBUG)
                    return None
                await asyncio.sleep(base_delay * (2 ** attempt))
            except Exception as e:
                self._log(f"  ✗ Oväntat hämtningsfel: {url} ({e})", LogLevel.DEBUG)
                return None
        return None

    async def _load_robots_txt(self):
        sitemaps_found = False
        try:
            r = await self.fetch(f"{self.base_url}/robots.txt", max_retries=1)
            if r and r.text:
                self.robot_parser = RobotFileParser()
                self.robot_parser.parse(r.text.splitlines())
                self._log("✓ robots.txt inläst")

                delay = self.robot_parser.crawl_delay("*")
                if delay:
                    self.rate_limiter.delay = float(delay)

                if self.find_sitemap:
                    sitemaps = [
                        line.split(': ', 1)[1].strip()
                        for line in r.text.splitlines()
                        if line.lower().startswith('sitemap:')
                    ]
                    if sitemaps:
                        await asyncio.gather(
                            *[self._parse_sitemap(sm) for sm in sitemaps],
                            return_exceptions=True
                        )
                        sitemaps_found = True
        except Exception as e:
            self._log(f"Kunde inte läsa robots.txt: {e}", LogLevel.DEBUG)

        if self.find_sitemap and not sitemaps_found:
            self._log("Letar efter sitemap.xml...")
            await self._parse_sitemap(f"{self.base_url}/sitemap.xml")

    async def _parse_sitemap(self, url: str):
        async with self.async_sitemap_lock:
            if url in self.visited_sitemaps:
                return
            self.visited_sitemaps.add(url)

        self._log(f"🗺️ Letar i sitemap: {url}")
        try:
            r = await self.fetch(url, max_retries=2, decode_text=False)
            if not r or not r.body:
                return
            content = r.body
            if url.lower().endswith('.gz'):
                content = gzip.decompress(content)

            try:
                soup = BeautifulSoup(content, 'lxml-xml')
                sitemap_urls = [
                    loc.text.strip() for sm in soup.find_all('sitemap')
                    if (loc := sm.find('loc'))
                ]
                url_strs = [
                    loc.text.strip() for node in soup.find_all('url')
                    if (loc := node.find('loc'))
                ]
            except Exception:
                import xml.etree.ElementTree as ET
                sitemap_urls = []
                url_strs = []
                try:
                    root = ET.fromstring(content)
                    for elem in root.iter():
                        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                        if tag == 'sitemap':
                            for child in elem:
                                ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                                if ctag == 'loc' and child.text:
                                    sitemap_urls.append(child.text.strip())
                        elif tag == 'url':
                            for child in elem:
                                ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                                if ctag == 'loc' and child.text:
                                    url_strs.append(child.text.strip())
                except Exception:
                    pass

            if sitemap_urls:
                await asyncio.gather(
                    *[self._parse_sitemap(s) for s in sitemap_urls],
                    return_exceptions=True
                )

            count = 0
            for s in url_strs:
                if self.is_valid_url(s):
                    # Spåra sitemap-funna dokument i manifestet — så att
                    # PDFer som bara hittas via sitemap (och inte via en
                    # webbsida) ändå syns där, om än utan riktig referer.
                    if (self.config.get("download_docs", False)
                            and self._looks_like_document_url(s)):
                        await self.manifest.record_link(
                            doc_url=s,
                            referer_url="",          # ingen riktig sida som länkar
                            referer_title="",        # → konvertern faller till metadata/filnamn
                            link_text="(via sitemap)",
                        )
                    if self.url_queue.add_url(s, depth=0,
                                              base_priority=CrawlPriority.SITEMAP.value):
                        count += 1
            if count > 0:
                self._log(f"✓ Hittade {count} (godkända) URLs i sitemap/index")
        except Exception as e:
            self._log(f"⚠ Fel vid sitemap-läsning: {e}", LogLevel.DEBUG)

    # ─── Playwright (lazy initialization) ──────────────────
    async def get_playwright_context(self):
        if not HAS_PLAYWRIGHT:
            return None
        if self._context is None:
            async with self.async_pw_lock:
                if self._context is None:
                    self._pw = await async_playwright().start()
                    is_headless = self.config.get("headless_mode", "headless") != "visible"
                    self._browser = await self._pw.chromium.launch(headless=is_headless)
                    self._context = await self._browser.new_context(ignore_https_errors=True)
                    if self.saved_cookies:
                        try:
                            await self._context.add_cookies(self.saved_cookies)
                        except Exception as e:
                            self._log(f"Kunde inte sätta Playwright-cookies: {e}",
                                      LogLevel.DEBUG)
        return self._context

    async def _render_with_playwright(self, url: str) -> Tuple[Optional[str], Optional[str]]:
        """Returnerar (html, final_url). Avbryter snabbt vid stop."""
        async with self.playwright_semaphore:
            ctx = await self.get_playwright_context()
            if ctx is None:
                return None, None
            page = await ctx.new_page()

            async def intercept_route(route):
                try:
                    if route.request.resource_type in (
                        "image", "stylesheet", "font", "media"
                    ):
                        await route.abort()
                    else:
                        await route.continue_()
                except Exception:
                    pass

            try:
                await page.route("**/*", intercept_route)
                # Kort timeout — om sidan inte är klar inom 12s tar vi vad vi har
                goto_ok = False
                try:
                    await page.goto(url, wait_until="domcontentloaded", timeout=12000)
                    goto_ok = True
                except Exception:
                    pass
                # Avbryt om goto misslyckades eller crawl stoppats
                if not goto_ok or self.state == CrawlerState.STOPPED:
                    return None, None
                html = await page.content()
                final_url = page.url
                return html, final_url
            finally:
                try:
                    await page.close()
                except Exception:
                    pass

    # ─── PII-tvätt ──────────────────────────────────────────
    def clean_pii(self, text: str) -> str:
        if not text:
            return text
        if self.config.get("remove_email"):
            text = re.sub(
                r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b',
                '[E-POST]', text
            )
        if self.config.get("remove_pnr"):
            pnr_pattern = (
                r'(?<!\d)(?:19|20)?\d{2}(?:0[1-9]|1[0-2])'
                r'(?:0[1-9]|[12]\d|3[01]|[6-9]\d)[\-\+]?\d{4}(?!\d)'
            )
            text = re.sub(pnr_pattern, '[PERSONNUMMER]', text)
        if self.config.get("remove_phone"):
            phone_pattern = (
                r'(?<!\d)(?:(?:\+|00)46[\s\-]*\(?0\)?[\s\-]*[1-9]|'
                r'0[\s\-]*\(?[1-9]\)?)[\s\-]*\d(?:[\s\-]*\d){4,8}\b'
            )
            text = re.sub(phone_pattern, '[TELEFON]', text)
        if self.config.get("remove_ip"):
            text = re.sub(r'\b(?:\d{1,3}\.){3}\d{1,3}\b', '[IP-ADRESS]', text)
        return text
    # ─── Innehållsextraktion ────────────────────────────────
    def extract_structured_data(self, html: str, url: str) -> Dict:
        try:
            soup = BeautifulSoup(html, 'lxml')
        except Exception:
            soup = BeautifulSoup(html, 'html.parser')

        page_links = []
        for a_tag in soup.find_all('a', href=True):
            href = a_tag.get('href', '').strip()
            if href and not href.startswith(('#', 'javascript:', 'mailto:')):
                full_url = urljoin(url, href)
                a_tag['href'] = full_url
                link_text = a_tag.get_text(separator=' ', strip=True)[:200]
                page_links.append({"url": full_url, "text": link_text})

        title = soup.title.string.strip() if soup.title and soup.title.string else "Okänd"
        title = self.clean_pii(title)

        keywords = []
        meta_kw = soup.find('meta', attrs={'name': re.compile(r'keywords', re.I)})
        if meta_kw and meta_kw.get('content'):
            keywords = [self.clean_pii(k.strip()) for k in meta_kw['content'].split(',')]

        description = ""
        meta_desc = soup.find('meta', attrs={'name': re.compile(r'description', re.I)})
        if meta_desc and meta_desc.get('content'):
            description = self.clean_pii(meta_desc['content'].strip())

        author_tag = soup.find('meta', attrs={'name': ['author', 'DC.creator']})
        author = self.clean_pii(author_tag['content']) if author_tag and author_tag.get('content') else ""

        pub_date = ""
        mod_date = ""

        meta_pub = soup.find('meta', attrs={'property': re.compile(
            r'article:published_time|og:pubdate', re.I)}) or \
            soup.find('meta', attrs={'name': re.compile(r'pubdate|date', re.I)})
        if meta_pub and meta_pub.get('content'):
            pub_date = meta_pub['content'].strip()

        meta_mod = soup.find('meta', attrs={'property': re.compile(
            r'article:modified_time|og:updated_time', re.I)}) or \
            soup.find('meta', attrs={'name': re.compile(r'last-modified|revised', re.I)}) or \
            soup.find('meta', attrs={'itemprop': 'dateModified'})
        if meta_mod and meta_mod.get('content'):
            mod_date = meta_mod['content'].strip()

        if not mod_date:
            time_tag = soup.find('time', attrs={'itemprop': 'dateModified'}) or \
                soup.find('time', class_=re.compile(r'update|modify', re.I))
            if time_tag:
                mod_date = time_tag.get('datetime', time_tag.get_text(strip=True))

        if not mod_date or not pub_date:
            for script in soup.find_all('script', type='application/ld+json'):
                if script.string:
                    if not mod_date:
                        m = re.search(r'"dateModified"\s*:\s*"([^"]+)"', script.string)
                        if m:
                            mod_date = m.group(1)
                    if not pub_date:
                        m = re.search(r'"datePublished"\s*:\s*"([^"]+)"', script.string)
                        if m:
                            pub_date = m.group(1)

        if not mod_date or not pub_date:
            text_content = soup.get_text(separator=' ', strip=True)
            date_pattern = (
                r'(?i)(?:senast\s+uppdaterad|uppdaterad|publicerad|ändrad)'
                r'[\s\:\*]*(?P<date>\d{1,2}\s+(?:januari|februari|mars|april|maj|juni|'
                r'juli|augusti|september|oktober|november|december)\s+\d{4}|\d{4}-\d{2}-\d{2})'
            )
            matches = list(re.finditer(date_pattern, text_content))
            if matches:
                extracted = matches[-1].group('date')
                if not mod_date and "uppdaterad" in matches[-1].group(0).lower():
                    mod_date = extracted
                if not pub_date and "publicerad" in matches[-1].group(0).lower():
                    pub_date = extracted
                if not mod_date and not pub_date:
                    mod_date = extracted

        lang_tag = soup.find('html')
        lang = lang_tag.get('lang', '') if lang_tag else ""

        og_type_tag = soup.find('meta', attrs={'property': 'og:type'})
        og_type = og_type_tag['content'] if og_type_tag and og_type_tag.get('content') else ""

        full_text = ""
        structured_sections: List[Dict] = []

        if self.use_trafilatura and HAS_TRAFILATURA:
            try:
                extracted = trafilatura.extract(
                    html, include_links=True, include_images=False,
                    include_tables=True, include_formatting=True,
                    output_format="markdown", url=url,
                )
                if extracted:
                    extracted = self.clean_pii(extracted)
                    extracted = absolutize_markdown_links(extracted, url)
                    raw_sections = re.split(r'(?=^#{1,6}\s)', extracted, flags=re.MULTILINE)
                    for raw in raw_sections:
                        raw = raw.strip()
                        if not raw:
                            continue
                        heading_match = re.match(r'^#{1,6}\s+(.+)', raw)
                        if heading_match:
                            heading = heading_match.group(1).strip()
                            content = raw[heading_match.end():].strip()
                        else:
                            heading = "Huvudinnehåll"
                            content = raw
                        if content:
                            structured_sections.append({"heading": heading, "text": content})
                    full_text = extracted
            except Exception as e:
                self._log(f"  ⚠ Trafilatura misslyckades ({e})", LogLevel.DEBUG)

        if not full_text:
            noise = re.compile(r'cookie|banner|menu|nav|sidebar|footer|share|social', re.I)
            for tag in ('script', 'style', 'nav', 'footer', 'aside',
                        'iframe', 'svg', 'button', 'form'):
                for el in soup.find_all(tag):
                    el.decompose()
            for el in soup.find_all(attrs={"class": noise}):
                el.decompose()
            for el in soup.find_all(attrs={"id": noise}):
                el.decompose()

            sections: List[Dict] = []
            current_heading = "Huvudinnehåll"
            current_text: List[str] = []

            for el in soup.find_all(['h1', 'h2', 'h3', 'h4', 'p', 'ul', 'ol', 'table']):
                if el.name.startswith('h'):
                    if current_text:
                        text_block = "\n".join(current_text).strip()
                        if text_block:
                            sections.append({"heading": current_heading, "text": text_block})
                    current_heading = el.get_text(separator=' ', strip=True)
                    current_text = []
                elif el.name in ('ul', 'ol'):
                    for li in el.find_all('li'):
                        parts = []
                        for child in li.children:
                            if hasattr(child, 'name') and child.name == 'a' and child.get('href'):
                                link_text = child.get_text(strip=True)
                                parts.append(f"[{link_text}]({child['href']})")
                            else:
                                t = child.get_text(strip=True) if hasattr(child, 'get_text') else str(child).strip()
                                if t:
                                    parts.append(t)
                        txt = " ".join(parts)
                        if txt:
                            current_text.append(f"• {txt}")
                elif el.name == 'table':
                    rows = el.find_all('tr')
                    for i, row in enumerate(rows):
                        cols = [c.get_text(separator=' ', strip=True)
                                for c in row.find_all(['td', 'th'])]
                        if cols:
                            current_text.append("| " + " | ".join(cols) + " |")
                            if i == 0:
                                current_text.append("|" + "|".join(["---"] * len(cols)) + "|")
                else:
                    parts = []
                    for child in el.children:
                        if hasattr(child, 'name') and child.name == 'a' and child.get('href'):
                            link_text = child.get_text(strip=True)
                            parts.append(f"[{link_text}]({child['href']})")
                        else:
                            t = child.get_text(strip=True) if hasattr(child, 'get_text') else str(child).strip()
                            if t:
                                parts.append(t)
                    txt = " ".join(parts)
                    if len(txt) > 5:
                        current_text.append(txt)

            if current_text:
                text_block = "\n".join(current_text).strip()
                if text_block:
                    sections.append({"heading": current_heading, "text": text_block})

            for s in sections:
                s['heading'] = self.clean_pii(s['heading'])
                s['text'] = self.clean_pii(s['text'])

            structured_sections = sections
            full_text = "\n\n".join([f"## {s['heading']}\n{s['text']}" for s in sections])

        return {
            "title": title,
            "url": url,
            "crawled_at": datetime.now().isoformat(),
            "author": author,
            "published_date": pub_date,
            "modified_date": mod_date,
            "language": lang,
            "og_type": og_type,
            "description": description,
            "keywords": keywords,
            "plain_text": full_text,
            "chunks": semantic_chunk_text(structured_sections, source_url=url),
            "page_links": page_links,
        }

    def _looks_like_document_url(self, url: str) -> bool:
        """Heuristik: är URL:en sannolikt ett binärt dokument?"""
        ext = posixpath.splitext(urlparse(url).path.lower())[1]
        return ext in DOCUMENT_EXTENSIONS

    def _needs_javascript(self, html: str) -> bool:
        """Avgör om sidan behöver Playwright-rendering.

        Kräver tydliga signaler. Tomma sidor triggar bara fallback om de
        uttryckligen ber om JavaScript eller innehållet är trivialt litet
        OCH det finns en SPA-rotnod.
        """
        if not html or len(html) < 500:
            return True
        html_lower = html.lower()
        if 'enable javascript' in html_lower or 'please enable javascript' in html_lower:
            return True
        # SPA-mönster: liten initial HTML + tom rotnod
        if len(html) < 2000:
            for pattern in ('id="root"', 'id="app"', 'id="__next"',
                            'ng-app', 'data-reactroot'):
                if pattern in html_lower:
                    return True
        return False

    # ─── Process page ───────────────────────────────────────
    async def process_page(self, url: str, depth: int) -> bool:
        if self.state == CrawlerState.STOPPED:
            return False
        while self.state == CrawlerState.PAUSED:
            await asyncio.sleep(0.3)

        domain = urlparse(url).netloc
        await self.rate_limiter.async_wait(domain)

        if self.state == CrawlerState.STOPPED:
            return False

        cached = await self.db.get_cache(url) if self.config.get("incremental") else None
        html, source, final_url = "", "Standard", url
        etag, last_mod = None, None

        try:
            # ─── Dokument-URL: hoppa över HTML-flödet helt ───
            # URL-extensionen är en starkare signal än Content-Type. Sitevision
            # och andra CMS:er returnerar ofta text/html för PDF-URL:er när
            # auth-cookies finns (de serverar en förhandsvyssida). Att lita på
            # HEAD/GET-content-type förlorade tidigare 1000+ dokument per crawl.
            if (self.config.get("download_docs", False)
                    and self._looks_like_document_url(url)):
                await self.download_document(url)
                return True

            # ─── HEAD bara för dokument-URL:er, inte HTML ───
            if self.use_hybrid and self._looks_like_document_url(url):
                head = await self.fetch(url, method='HEAD', max_retries=1)
                if head and head.content_type:
                    ct = head.content_type
                    if any(t in ct for t in ('image/', 'video/', 'audio/', 'font/')):
                        return False
                    doc_types = ('application/pdf', 'application/vnd', 'application/msword')
                    if any(dt in ct for dt in doc_types):
                        if self.config.get("download_docs", False):
                            await self.download_document(url)
                        return True

            # ─── GET (med conditional headers) ───
            if self.use_hybrid:
                result = await self.fetch(url, cached=cached)
                if result is None:
                    async with self.async_stats_lock:
                        self.stats.pages_failed += 1
                    self._gui_update(url, "Fel", "")
                    return False

                # 304 Not Modified — uppdatera bara timestamp och gå vidare
                if result.not_modified:
                    async with self.async_stats_lock:
                        self.stats.pages_not_modified_304 += 1
                        self.stats.pages_visited += 1
                    self._gui_update(url, "Ej ändrad (304)",
                                     cached.get('title', '') if cached else '')
                    await self.db.touch_cache(url, etag=result.etag,
                                              last_modified=result.last_modified)
                    # Hoppa över extraktion och länkdetektering — vi har redan dessa länkar
                    return True

                # 401/403 → behandla som login-utgång eller räkna som fel
                if result.status in (401, 403):
                    if self._login_detection_enabled:
                        await self._handle_login_expired(url)
                    else:
                        self._log(f"  ✗ HTTP {result.status}: {url}", LogLevel.DEBUG)
                        async with self.async_stats_lock:
                            self.stats.pages_failed += 1
                        self._gui_update(url, f"HTTP {result.status}", "")
                    return False

                # Content-typ-hantering
                final_url = result.final_url or url
                etag = result.etag
                last_mod = result.last_modified
                content_type = result.content_type

                if any(t in content_type for t in ('image/', 'video/', 'audio/', 'font/')):
                    return False
                if any(dt in content_type for dt in ('application/pdf', 'application/vnd',
                                                     'application/msword')):
                    if self.config.get("download_docs", False):
                        await self.download_document(url)
                    return True

                html = result.text or ""
                if not html:
                    return False

                # JS-fallback
                if self._needs_javascript(html):
                    if self.state == CrawlerState.STOPPED:
                        return False
                    rendered_html, rendered_url = await self._render_with_playwright(url)
                    if rendered_html:
                        html = rendered_html
                        final_url = rendered_url or final_url
                        source = "Webbläsare"
                        # ETag/Last-Modified från aiohttp gäller JS-skalet,
                        # inte det renderade innehållet — nollställ så att
                        # nästa crawl inte missar ändringar via felaktig 304.
                        etag = None
                        last_mod = None
                        async with self.async_stats_lock:
                            self.stats.playwright_fallbacks += 1
            else:
                # Endast Playwright
                if self.state == CrawlerState.STOPPED:
                    return False
                rendered_html, rendered_url = await self._render_with_playwright(url)
                if rendered_html:
                    html = rendered_html
                    final_url = rendered_url or url
                    source = "Webbläsare"

            if self.state == CrawlerState.STOPPED:
                return False

            # ─── Login-detektion ───
            if self._login_detection_enabled and html:
                # 'final_url' kan vara samma som 'url' — då är det ingen redirect.
                # Vi kollar bara om sidan redirectades till login-URL ELLER om
                # innehållet otvetydigt är ett login-formulär.
                redirected_away = (normalize_url(final_url) != normalize_url(url))
                triggered = False
                if redirected_away and self.login_detector.is_login_redirect(url, final_url):
                    triggered = True
                elif self.login_detector.is_login_content(html):
                    triggered = True

                if triggered:
                    await self._handle_login_expired(url)
                    return False

            data = await asyncio.to_thread(self.extract_structured_data, html, url)
            content_hash = get_clean_hash(data["plain_text"])
            text_length = len(data["plain_text"])

            if self.config.get("incremental") and cached and cached.get('hash') == content_hash:
                async with self.async_stats_lock:
                    self.stats.pages_unchanged += 1
                self._gui_update(url, "Oförändrad", data["title"])
            else:
                self._gui_update(url, f"Hämtad ({source})", data["title"])

                if text_length > 50 and self.save_format not in ("Ingen text", "No text"):
                    texts_dir = os.path.join(self.output_dir, "texter")
                    os.makedirs(texts_dir, exist_ok=True)
                    fn = stable_filename(url, self.save_format)
                    out_path = os.path.join(texts_dir, fn)

                    if self.save_format == ".json":
                        # Behåll plain_text i JSON-output (chunks finns redan separat)
                        with open(out_path, 'w', encoding='utf-8') as f:
                            json.dump(data, f, ensure_ascii=False, indent=2)
                    else:
                        # Markdown: bädda in URL i BRÖDTEXTEN (inte bara i header)
                        # eftersom RAG-pipelines ofta kapar de första raderna före
                        # retrieval. Upprepa källan sist så att även avslutande
                        # chunks har källinformation.
                        with open(out_path, 'w', encoding='utf-8') as f:
                            body = data['plain_text']
                            body = strip_cms_boilerplate(body)
                            body = downgrade_body_h1(body)
                            f.write(
                                f"# {data['title']}\n\n"
                                f"**Källa:** {url}\n\n"
                                f"---\n\n"
                                f"{body}"
                                f"\n\n---\n\n"
                                f"**Källa:** {url}\n"
                            )

            await self.db.save_cache(url, content_hash, data["title"], text_length,
                                     etag=etag, last_modified=last_mod)

            # ─── Länkdetektering ───
            if self.max_depth == 0 or depth < self.max_depth:
                page_title = data.get("title", "")
                for link_info in data.get("page_links", []):
                    full_url = link_info["url"]
                    if not self.is_valid_url(full_url):
                        continue

                    # Om länken pekar på ett dokument: registrera referer-info
                    # i manifestet INNAN länken hamnar i kön. Vi har bara
                    # tillgång till sidans titel och länktexten just nu —
                    # dokument-URL:en själv har ingen sidkontext.
                    if (self.config.get("download_docs", False)
                            and self._looks_like_document_url(full_url)):
                        await self.manifest.record_link(
                            doc_url=full_url,
                            referer_url=url,
                            referer_title=page_title,
                            link_text=link_info["text"],
                        )

                    self.url_queue.add_url(full_url, depth=depth + 1)

            async with self.async_stats_lock:
                self.stats.pages_visited += 1
            return True

        except asyncio.CancelledError:
            raise
        except Exception as e:
            self._log(f"  ✗ Fel vid besök ({url}): {str(e)[:80]}", LogLevel.ERROR)
            async with self.async_stats_lock:
                self.stats.pages_failed += 1
            self._gui_update(url, "Fel", str(e)[:30])
            return False

    async def _handle_login_expired(self, url: str):
        self._log(f"⚠ Session utgången — inloggningssida detekterad för: {url}",
                  LogLevel.WARNING)
        async with self.async_stats_lock:
            self.stats.pages_failed += 1
        self._gui_update(url, "Session utgången", "")

    async def download_document(self, url: str):
        if self.state == CrawlerState.STOPPED:
            return
        while self.state == CrawlerState.PAUSED:
            await asyncio.sleep(0.3)

        async with self.async_download_lock:
            if url in self.downloaded_files:
                return
            self.downloaded_files.add(url)

        # Rate limit även dokumentnedladdningar
        domain = urlparse(url).netloc
        await self.rate_limiter.async_wait(domain)

        try:
            docs_dir = os.path.join(self.output_dir, "dokument")
            os.makedirs(docs_dir, exist_ok=True)

            # Pre-compute filnamn från URL (stabilt, behövs för Playwright-fallback)
            url_basename = os.path.basename(unquote(urlparse(url).path)) or ""
            url_ext = posixpath.splitext(url_basename)[1].lower() if url_basename else ""
            # Fallback-extension bestäms nedan av Content-Type/Content-Disposition
            # om URL:en saknar ändelse — ".bin" som säker default
            if not url_ext:
                url_ext = ".bin"
            slug_base = slugify((url_basename.split('.')[0] if url_basename else 'dokument'))[:50]
            url_hash = hashlib.md5(url.encode('utf-8')).hexdigest()[:6]
            safe_filename = f"{slug_base}_{url_hash}{url_ext}"
            filepath = os.path.join(docs_dir, safe_filename)

            # ─── Filen finns redan: registrera bara i manifestet ───
            if os.path.exists(filepath):
                bytes_written = os.path.getsize(filepath)
                await self.manifest.record_download(
                    doc_url=url, filename=safe_filename,
                    size_bytes=bytes_written,
                )
                if self.converter and self.converter.can_convert(filepath):
                    base = slugify(os.path.splitext(safe_filename)[0])[:50]
                    url_hash = hashlib.md5(url.encode('utf-8')).hexdigest()[:6]
                    md_check = os.path.join(
                        self.converter.texts_dir,
                        f"{base}_{url_hash}_doc.md")
                    if not os.path.exists(md_check):
                        referers = self.manifest._referers.get(url, [])
                        ref = referers[0] if referers else {}
                        md_path = await asyncio.to_thread(
                            self.converter.convert,
                            filepath, url,
                            referer_url=ref.get('referer_url', ''),
                            referer_title=ref.get('referer_title', ''),
                            link_text=ref.get('link_text', ''),
                        )
                        if md_path:
                            self._log(
                                f"  📄 Konverterad till .md: "
                                f"{os.path.basename(md_path)}",
                                LogLevel.DEBUG)
                self._log(f"  ↩ Fanns redan, registrerad i manifest: {safe_filename}",
                          LogLevel.DEBUG)
                return

            # ─── Försök 1: aiohttp ───
            aiohttp_ok = False
            try:
                async with self.req_session.get(
                    url, timeout=aiohttp.ClientTimeout(total=60)
                ) as resp:
                    if resp.status != 200:
                        self._log(
                            f"  ⚠ aiohttp HTTP {resp.status} för dokument: {url}",
                            LogLevel.DEBUG)
                    else:
                        content_type = resp.headers.get('Content-Type', '').lower()
                        cd_header = resp.headers.get('Content-Disposition', '').lower()
                        is_attachment = 'attachment' in cd_header
                        doc_content_types = (
                            'application/pdf', 'application/vnd',
                            'application/msword', 'application/octet-stream',
                            'application/x-download', 'application/force-download',
                            'application/zip', 'application/x-zip',
                            'application/x-rar',
                        )
                        looks_like_doc = (
                            is_attachment
                            or any(dt in content_type for dt in doc_content_types)
                        )
                        if not looks_like_doc and 'text/html' in content_type:
                            self._log(
                                f"  ⚠ Servern svarar HTML istället för dokument: {url}",
                                LogLevel.DEBUG)
                        else:
                            # Härleda filändelse från Content-Type om URL:en
                            # saknade extension (url_ext == ".bin")
                            if url_ext == ".bin":
                                ct_ext_map = {
                                    'application/pdf': '.pdf',
                                    'application/msword': '.doc',
                                    'application/vnd.openxmlformats-officedocument.wordprocessingml': '.docx',
                                    'application/vnd.openxmlformats-officedocument.spreadsheetml': '.xlsx',
                                    'application/vnd.openxmlformats-officedocument.presentationml': '.pptx',
                                    'application/vnd.ms-excel': '.xls',
                                    'application/vnd.ms-powerpoint': '.ppt',
                                    'application/zip': '.zip',
                                    'application/x-zip': '.zip',
                                }
                                for ct_prefix, ct_extension in ct_ext_map.items():
                                    if ct_prefix in content_type:
                                        url_ext = ct_extension
                                        safe_filename = f"{slug_base}_{url_hash}{url_ext}"
                                        filepath = os.path.join(docs_dir, safe_filename)
                                        break

                            # Uppdatera filnamn från Content-Disposition om tillgängligt
                            cd_raw = resp.headers.get('Content-Disposition', '')
                            if cd_raw:
                                match = re.search(
                                    r"filename\*?=(?:UTF-8'')?[\"']?([^\"';]+)[\"']?",
                                    cd_raw, flags=re.IGNORECASE)
                                if match:
                                    cd_name = unquote(match.group(1)).strip()
                                    if cd_name:
                                        cd_ext = posixpath.splitext(cd_name)[1].lower() or url_ext
                                        safe_filename = (
                                            f"{slugify(cd_name.split('.')[0])[:50]}_"
                                            f"{hashlib.md5(url.encode('utf-8')).hexdigest()[:6]}"
                                            f"{cd_ext}"
                                        )
                                        filepath = os.path.join(docs_dir, safe_filename)

                            # Ladda ner
                            bytes_written = 0
                            aborted = False
                            with open(filepath, 'wb') as f:
                                while True:
                                    if self.state == CrawlerState.STOPPED:
                                        aborted = True
                                        break
                                    chunk = await resp.content.read(8192)
                                    if not chunk:
                                        break
                                    f.write(chunk)
                                    bytes_written += len(chunk)

                            if aborted:
                                try:
                                    os.remove(filepath)
                                except OSError:
                                    pass
                                async with self.async_download_lock:
                                    self.downloaded_files.discard(url)
                                return

                            aiohttp_ok = True
            except (aiohttp.ClientError, asyncio.TimeoutError) as e:
                self._log(f"  ⚠ aiohttp-fel vid dokumentnedladdning: {e}",
                          LogLevel.DEBUG)

            # ─── Försök 2: Playwright-fallback (SAML-skyddade dokument) ───
            if not aiohttp_ok and HAS_PLAYWRIGHT:
                self._log(f"  🔄 Försöker Playwright-fallback för: {url}",
                          LogLevel.DEBUG)
                aiohttp_ok = await self._download_via_playwright(url, filepath)

            if not aiohttp_ok:
                self._log(f"  ✗ Kunde inte ladda ner dokument: {url}",
                          LogLevel.WARNING)
                async with self.async_download_lock:
                    self.downloaded_files.discard(url)
                # Rensa ev. tom/korrupt fil
                if os.path.exists(filepath) and os.path.getsize(filepath) == 0:
                    try:
                        os.remove(filepath)
                    except OSError:
                        pass
                return

            # ─── Framgång: stats, manifest, konvertering ───
            bytes_written = os.path.getsize(filepath)
            async with self.async_stats_lock:
                self.stats.documents_downloaded += 1
                self.stats.bytes_downloaded += bytes_written
            await self.manifest.record_download(
                doc_url=url, filename=safe_filename,
                size_bytes=bytes_written,
            )
            self._log(f"  ⬇ Dokument sparat: {safe_filename}")

            # Konvertera till Markdown om aktiverat
            if self.converter and self.converter.can_convert(filepath):
                referers = self.manifest._referers.get(url, [])
                ref = referers[0] if referers else {}
                md_path = await asyncio.to_thread(
                    self.converter.convert,
                    filepath, url,
                    referer_url=ref.get('referer_url', ''),
                    referer_title=ref.get('referer_title', ''),
                    link_text=ref.get('link_text', ''),
                )
                if md_path:
                    self._log(
                        f"  📄 Konverterad till .md: "
                        f"{os.path.basename(md_path)}")
        except Exception as e:
            self._log(f"  ✗ Filnedladdning misslyckades: {url}, {e}", LogLevel.ERROR)
            async with self.async_download_lock:
                self.downloaded_files.discard(url)

    async def _download_via_playwright(self, url: str, filepath: str) -> bool:
        """Fallback-nedladdning via Playwright för SAML-skyddade dokument.

        Anropas när aiohttp misslyckas (HTTP-fel eller SAML-redirect till
        login-sida). Playwright har den fulla browser-sessionen med
        SAML-cookies och klarar de redirect-kedjor som aiohttp missar.

        Två strategier:
          1. expect_download — servern skickar Content-Disposition: attachment
          2. context.request.get — inline-dokument (PDF i browsern)
        """
        async with self.playwright_semaphore:
            ctx = await self.get_playwright_context()
            if ctx is None:
                return False

            page = await ctx.new_page()
            try:
                # Strategi 1: Fånga en download-händelse
                try:
                    async with page.expect_download(timeout=15000) as dl_info:
                        await page.goto(url, wait_until="commit", timeout=15000)
                    download = await dl_info.value
                    await download.save_as(filepath)
                    self._log(
                        f"  🔄 Playwright-download lyckades: "
                        f"{os.path.basename(filepath)}", LogLevel.DEBUG)
                    return True
                except Exception:
                    pass

                # Strategi 2: API-request med browser-cookies
                try:
                    api_resp = await ctx.request.get(url, timeout=20000)
                    if api_resp.ok:
                        ct = (api_resp.headers.get('content-type') or '').lower()
                        if 'text/html' not in ct:
                            body = await api_resp.body()
                            if body and len(body) > 100:
                                with open(filepath, 'wb') as f:
                                    f.write(body)
                                self._log(
                                    f"  🔄 Playwright-API lyckades: "
                                    f"{os.path.basename(filepath)}",
                                    LogLevel.DEBUG)
                                return True
                except Exception:
                    pass

                return False
            finally:
                try:
                    await page.close()
                except Exception:
                    pass


    async def _generate_index(self):
        self._log("📊 Skapar index-fil (index.csv)...")
        try:
            records = await self.db.get_all_records()
            with open(os.path.join(self.output_dir, "index.csv"),
                      'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                writer.writerow(['URL', 'Titel', 'Hämtad_Datum', 'Filnamn'])
                for url, title, date, content_hash in records:
                    fn = stable_filename(url, self.save_format)
                    writer.writerow([url, title, date, fn])
        except Exception as e:
            self._log(f"⚠ Fel vid skapande av index.csv: {e}", LogLevel.ERROR)

    async def _generate_manifest(self):
        """Skriver dokument-manifest.json till utmappen.

        Filen länkar varje nedladdat dokument till den sida som hade
        länken — RAG-systemet kan slå upp ett filnamn där och få
        tillbaka rätt intranät-URL att citera som källa.
        """
        try:
            manifest_data = self.manifest.build(domain=self.domain)
            count = manifest_data.get("document_count", 0)
            if count == 0:
                # Inga dokument laddades ner — ingen anledning att skapa manifest
                return
            manifest_filename = f"manifest_{slugify(self.domain)}.json"
            manifest_path = os.path.join(self.output_dir, manifest_filename)
            with open(manifest_path, 'w', encoding='utf-8') as f:
                json.dump(manifest_data, f, ensure_ascii=False, indent=2)
            orphans = manifest_data.get("orphan_count", 0)
            msg = f"📋 Manifest skapad: {count} dokument"
            if orphans:
                msg += f" ({orphans} länkade men ej nedladdade)"
            self._log(msg)
        except Exception as e:
            self._log(f"⚠ Fel vid skapande av manifest-fil: {e}", LogLevel.ERROR)

    def pause(self):
        if self.state == CrawlerState.RUNNING:
            self.state = CrawlerState.PAUSED
            return True
        elif self.state == CrawlerState.PAUSED:
            self.state = CrawlerState.RUNNING
            return False
        return False

    def stop(self):
        self.state = CrawlerState.STOPPED
        self._log("🛑 Avbryter crawl (väntar på aktiva processer)...")
        self.login_event.set()

    # ─── Huvudloopen ────────────────────────────────────────
    async def crawl(self):
        await self.db.connect()

        # Login-läge: starta synlig browser, vänta på user, plocka cookies
        if self.config["headless_mode"] == "login_then_headless":
            if HAS_PLAYWRIGHT:
                async with async_playwright() as p:
                    browser = await p.chromium.launch(headless=False)
                    context = await browser.new_context()
                    page = await context.new_page()

                    self._log("👤 Navigerar till start-URL för manuell inloggning...")
                    await page.goto(self.start_url)

                    self._log("\n⏳ VÄNTAR PÅ MANUELL INLOGGNING...")
                    if self.msg_queue:
                        self.msg_queue.put(("login_wait", None))
                    await asyncio.to_thread(self.login_event.wait)

                    if self.state == CrawlerState.STOPPED:
                        return

                    self._log("🔄 Sparar cookies och byter till osynligt läge...")
                    self.saved_cookies = await context.cookies()
                    await browser.close()
            else:
                self._log("⚠ Playwright saknas, kan inte utföra manuell inloggning!",
                          LogLevel.ERROR)

        self.req_session = await self._create_session()

        if self.config.get("respect_robots", True):
            await self._load_robots_txt()

        self.stats.start_time = datetime.now()

        try:
            self.active_tasks = 0

            async def bounded_process(url, depth):
                try:
                    while self.state == CrawlerState.PAUSED:
                        await asyncio.sleep(0.3)
                    if self.state == CrawlerState.STOPPED:
                        return

                    async with self.semaphore:
                        while self.state == CrawlerState.PAUSED:
                            await asyncio.sleep(0.3)
                        if self.state == CrawlerState.STOPPED:
                            return
                        await self.process_page(url, depth)
                except asyncio.CancelledError:
                    pass
                except Exception as e:
                    self._log(f"💥 Oväntat fel i bounded_process ({url}): {e}",
                              LogLevel.ERROR)
                finally:
                    try:
                        async with self.async_stats_lock:
                            self.active_tasks -= 1
                    except Exception:
                        self.active_tasks = max(0, self.active_tasks - 1)

            async with asyncio.TaskGroup() as tg:
                while self.state != CrawlerState.STOPPED:
                    if self.state == CrawlerState.PAUSED:
                        await asyncio.sleep(0.5)
                        continue
                    if self.max_pages > 0 and self.stats.pages_visited >= self.max_pages:
                        break

                    queue_item = self.url_queue.get_next()
                    if not queue_item:
                        async with self.async_stats_lock:
                            tasks_are_zero = (self.active_tasks == 0)
                        if self.url_queue.size() == 0 and tasks_are_zero:
                            break
                        await asyncio.sleep(0.3)
                        continue

                    async with self.async_stats_lock:
                        self.active_tasks += 1
                    tg.create_task(bounded_process(queue_item[1], queue_item[0]))

        except Exception as e:
            self._log(f"💥 Oväntat fel i async crawl: {e}", LogLevel.ERROR)
        finally:
            try:
                await self._generate_index()
            except Exception:
                pass
            try:
                await self._generate_manifest()
            except Exception:
                pass
            if self.req_session:
                await self.req_session.close()
            await self.db.close()
            if self._browser:
                try:
                    await self._browser.close()
                except Exception:
                    pass
            if self._pw:
                try:
                    await self._pw.stop()
                except Exception:
                    pass
            self.stats.end_time = datetime.now()
            rate = self.stats.pages_per_second
            self._log(f"Färdig! Total tid: {self.stats.duration} "
                      f"({rate:.1f} sidor/sek, "
                      f"{self.stats.pages_not_modified_304} via 304-cache)")
            if self.msg_queue:
                self.msg_queue.put(("done", "Klar"))


# ─────────────────────────────────────────────────────────────
#  SERVER / CLI LÄGE
# ─────────────────────────────────────────────────────────────
async def run_cli_mode(config_file: str, webhook_url: Optional[str] = None):
    print(f"🚀 Startar Webbdammsugare Pro Serverläge med filen: {config_file}")
    with open(config_file, 'r', encoding='utf-8') as f:
        sites_config = json.load(f)

    base_out = os.path.abspath("server_data")
    site_semaphore = asyncio.Semaphore(3)

    async def run_single_site(site):
        async with site_semaphore:
            site_out = os.path.join(base_out, slugify(site.get("name", "Unknown")))
            config = {
                **site,
                "output_dir": site_out,
                "headless_mode": site.get("headless_mode", "headless"),
                "use_hybrid": True,
                "incremental": True,
            }
            crawler = AsyncWebCrawler(config)
            await crawler.crawl()
            return site.get("name"), crawler

    tasks = [run_single_site(site) for site in sites_config]
    start_time = time.time()
    results = await asyncio.gather(*tasks, return_exceptions=True)

    summary_lines = []
    for r in results:
        if isinstance(r, Exception):
            summary_lines.append(f" - FEL: {r}")
        else:
            name, c = r
            summary_lines.append(
                f" - {name}: {c.stats.pages_visited} besökta, "
                f"{c.stats.pages_failed} fel, "
                f"{c.stats.pages_not_modified_304} via 304."
            )
    summary = "\n".join(summary_lines)
    print(f"\n✅ Klart på {time.time() - start_time:.1f} sekunder!\n{summary}")

    if webhook_url:
        try:
            requests.post(webhook_url,
                          json={"text": f"🚀 Nattens dammsugning klar!\n{summary}"})
        except Exception:
            pass
# ─────────────────────────────────────────────────────────────
#  GRAFISKT GRÄNSSNITT (GUI - CustomTkinter)
# ─────────────────────────────────────────────────────────────
class AppGUI:
    def __init__(self, root: ctk.CTk):
        self.root = root
        self.lang = "sv"
        self.texts = {
            "window_title": {"sv": "Webbdammsugare Pro (v7.0)", "en": "Web Crawler Pro (v7.0)"},
            "tab_basic": {"sv": "⚙️ Grundinställningar", "en": "⚙️ Basic Settings"},
            "tab_adv": {"sv": "🔧 Avancerat", "en": "🔧 Advanced"},
            "lbl_url": {"sv": "🌐 Startadress:", "en": "🌐 Start URL:"},
            "btn_help": {"sv": "❓ Hjälp", "en": "❓ Help"},
            "lbl_delay": {"sv": "Fördröjning (sek):", "en": "Delay (sec):"},
            "lbl_max_pages": {"sv": "Max sidor (0=Oändligt):", "en": "Max pages (0=Infinite):"},
            "lbl_max_depth": {"sv": "Max djup (0=Oändligt):", "en": "Max depth (0=Infinite):"},
            "lbl_format": {"sv": "Filformat:", "en": "File Format:"},
            "lbl_concurrency": {"sv": "Samtidighet:", "en": "Concurrency:"},
            "cb_docs": {"sv": "Ladda ner dokument (PDF m.m.)", "en": "Download documents (PDF etc.)"},
            "cb_convert_docs": {"sv": "📄 Konvertera dokument till Markdown", "en": "📄 Convert documents to Markdown"},
            "lbl_mode": {"sv": "Körläge:", "en": "Run Mode:"},
            "lbl_folder": {"sv": "Spara Mapp:", "en": "Save Folder:"},
            "btn_folder": {"sv": "Välj Mapp...", "en": "Browse..."},
            "cb_hybrid": {"sv": "⚡ Hybrid-motor (Requests + Playwright)", "en": "⚡ Hybrid Engine (Requests + Playwright)"},
            "cb_traf": {"sv": "🧠 Använd Trafilatura för text", "en": "🧠 Use Trafilatura for text extraction"},
            "cb_sitemap": {"sv": "Läs Sitemap.xml", "en": "Parse Sitemap.xml"},
            "cb_robots": {"sv": "Respektera robots.txt", "en": "Respect robots.txt"},
            "cb_strict": {"sv": "Strikt Domän", "en": "Strict Domain"},
            "lbl_exclude": {"sv": "Uteslut ord i URL:", "en": "Exclude words in URL:"},
            "lbl_require": {"sv": "Kräv ord i URL (något av):", "en": "Require words in URL (any of):"},
            "cb_rm_email": {"sv": "Radera E-post", "en": "Remove Email"},
            "cb_rm_phone": {"sv": "Radera Telefonnummer", "en": "Remove Phone Numbers"},
            "cb_rm_pnr": {"sv": "Radera Personnummer", "en": "Remove Swedish SSN"},
            "cb_rm_ip": {"sv": "Radera IP-adresser", "en": "Remove IP Addresses"},
            "btn_start": {"sv": "▶ Starta", "en": "▶ Start"},
            "btn_pause": {"sv": "⏸ Pausa", "en": "⏸ Pause"},
            "lbl_template": {"sv": "📋 Mall:", "en": "📋 Template:"},
            "template_none": {"sv": "— Ingen mall —", "en": "— No template —"},
            "template_loaded": {"sv": "✓ Mall laddad: {}", "en": "✓ Template loaded: {}"},
            "template_not_found": {"sv": "Ingen sites.json hittades", "en": "No sites.json found"},
            "btn_resume": {"sv": "▶ Fortsätt", "en": "▶ Resume"},
            "btn_stop": {"sv": "■ Stoppa", "en": "■ Stop"},
            "col_status": {"sv": "Status", "en": "Status"},
            "col_title": {"sv": "Sido-titel", "en": "Page Title"},
            "status_wait": {"sv": "Väntar på start...", "en": "Waiting to start..."},
            "stats_fmt": {
                "sv": "Besökta: {} | Sidor: {} | Dokument: {} | Cachade: {} | I Kö: {} | Fel: {} | Tid kvar: {}",
                "en": "Visited: {} | Pages: {} | Docs: {} | Cached: {} | Queued: {} | Errors: {} | ETA: {}"
            },
            "help_title": {"sv": "❓ Hjälp & Instruktioner", "en": "❓ Help & Instructions"},
            "run_modes": {
                "headless": {"sv": "Snabb (dold)", "en": "Fast (hidden)"},
                "login_then_headless": {"sv": "Logga in, sen dold", "en": "Login, then hidden"},
                "visible": {"sv": "Synlig (felsökning)", "en": "Visible (debugging)"}
            },
            "help_content": {
                "sv": ("⚙️ GRUNDINSTÄLLNINGAR\n-------------------------\n"
                       "* Startadress: URL där programmet börjar leta.\n"
                       "* Fördröjning: Tid mellan sidbesök (per domän).\n"
                       "* Max sidor/djup: 0 betyder oändligt.\n"
                       "* Samtidighet: antal parallella sidor (default 10).\n"
                       "* Körläge:\n"
                       "  - Snabb (dold): Snabbast, körs i bakgrunden.\n"
                       "  - Logga in: Öppnar fönster för inloggning, kör sen dolt.\n"
                       "* Filformat:\n"
                       "  - .json: Strukturerad data anpassad för Vektordatabaser och AI.\n"
                       "  - .md: Markdown, bra för generella LLM-läsningar.\n"
                       "  - Ingen text: Skrapar enbart dokument (om ikryssat).\n\n"
                       "📋 MALLAR\n-------------------------\n"
                       "Lägg en sites.json i samma mapp som programmet.\n"
                       "Välj en mall i dropdown-menyn så fylls alla inställningar i automatiskt.\n\n"
                       "🔧 AVANCERAT\n-------------------------\n"
                       "* Hybrid-motor: Rekommenderas för modern webb.\n"
                       "* URL-Filter: Filtrerar på ord i URL:en, inte i sidans text.\n"
                       "* PII-Tvätt: Raderar personuppgifter automatiskt innan sparning.\n\n"
                       "♻️ INKREMENTELL CRAWL\n-------------------------\n"
                       "Programmet kommer ihåg ETag/Last-Modified för varje sida.\n"
                       "Vid omstart svarar servern '304 Not Modified' för oförändrade sidor\n"
                       "— ofta 10-50× snabbare än en första körning.\n\n"
                       "📋 DOKUMENT-MANIFEST\n-------------------------\n"
                       "När du laddar ner dokument (PDF m.m.) skapas också en manifest.json\n"
                       "i utmappen som listar varje dokument tillsammans med vilken intranät-\n"
                       "sida som hade länken till det.\n\n"
                       "📄 DOKUMENT → MARKDOWN\n-------------------------\n"
                       "Kryssa i 'Konvertera dokument till Markdown' så extraheras texten\n"
                       "ur PDF, Word, Excel och PowerPoint och sparas som .md-filer med\n"
                       "intranät-URL:en i toppen. Perfekt för RAG-system som Svea där\n"
                       "AI:n behöver se källan direkt i texten för att citera rätt.\n"
                       "Kräver: pip install PyMuPDF python-docx openpyxl python-pptx\n\n"
                       "💻 SERVER-LÄGE\n-------------------------\n"
                       "Körs via CMD för automatisering:\n"
                       "python uwc.py --config sites.json\n\n"
                       "💡 TIPS: Dubbelklicka på en rad i tabellen för att öppna länken!"),
                "en": ("⚙️ BASIC SETTINGS\n-------------------------\n"
                       "* Start URL: Where the crawler begins.\n"
                       "* Delay: Seconds to wait between requests (per domain).\n"
                       "* Concurrency: number of parallel pages (default 10).\n"
                       "* Run Mode:\n"
                       "  - Fast (hidden): Fastest option, background.\n"
                       "  - Login: Shows browser for login, then background.\n"
                       "* File Format:\n"
                       "  - .json: Structured output for Vector Databases and AI.\n"
                       "  - .md: Markdown.\n"
                       "  - No text: Only downloads documents (if checked).\n\n"
                       "📋 TEMPLATES\n-------------------------\n"
                       "Place a sites.json in the same folder as the program.\n\n"
                       "🔧 ADVANCED\n-------------------------\n"
                       "* Hybrid Engine: Recommended for modern web.\n"
                       "* URL Filters: Filters on URL substrings.\n"
                       "* PII Wash: Removes personal data before saving.\n\n"
                       "♻️ INCREMENTAL CRAWL\n-------------------------\n"
                       "ETag and Last-Modified are cached per URL.\n"
                       "On re-runs, the server responds '304 Not Modified' for unchanged pages\n"
                       "— often 10-50x faster than the first crawl.\n\n"
                       "📋 DOCUMENT MANIFEST\n-------------------------\n"
                       "When documents (PDFs etc.) are downloaded, a manifest.json is also\n"
                       "created in the output folder.\n\n"
                       "📄 DOCUMENTS → MARKDOWN\n-------------------------\n"
                       "Check 'Convert documents to Markdown' to extract text from PDF,\n"
                       "Word, Excel and PowerPoint files and save as .md with the intranet\n"
                       "URL at the top. Perfect for RAG systems where the AI needs the\n"
                       "source URL directly in the text to cite correctly.\n"
                       "Requires: pip install PyMuPDF python-docx openpyxl python-pptx\n\n"
                       "💻 SERVER MODE\n-------------------------\n"
                       "python uwc.py --config sites.json\n\n"
                       "💡 TIP: Double-click a row to open the URL!")
            }
        }

        self.root.title(self.texts["window_title"][self.lang])
        self.root.geometry("1000x880")
        self.msg_queue = queue.Queue()
        self.crawler_instance: Optional[AsyncWebCrawler] = None
        self._update_treeview_style("Light")
        self._build_ui()
        self.root.protocol("WM_DELETE_WINDOW", self._on_closing)
        self.root.after(100, self.process_queue)

    def _on_closing(self):
        if self.crawler_instance:
            self.crawler_instance.stop()
        self.root.destroy()

    def change_appearance_mode_event(self):
        if self.theme_switch.get() == 1:
            ctk.set_appearance_mode("Dark")
            self.theme_switch.configure(text="🌙")
            self._update_treeview_style("Dark")
        else:
            ctk.set_appearance_mode("Light")
            self.theme_switch.configure(text="☀️")
            self._update_treeview_style("Light")

    def change_language_event(self, choice):
        inverted_map = {v[self.lang]: k for k, v in self.texts["run_modes"].items()}
        internal_mode = inverted_map.get(self.headless_var.get(), "headless")

        new_lang = "sv" if "SV" in choice else "en"
        if new_lang == self.lang:
            return

        old_basic = self.texts["tab_basic"][self.lang]
        new_basic = self.texts["tab_basic"][new_lang]
        old_adv = self.texts["tab_adv"][self.lang]
        new_adv = self.texts["tab_adv"][new_lang]
        self.tabview.rename(old_basic, new_basic)
        self.tabview.rename(old_adv, new_adv)

        self.lang = new_lang
        self.root.title(self.texts["window_title"][self.lang])

        self.lbl_url.configure(text=self.texts["lbl_url"][self.lang])
        self.help_btn.configure(text=self.texts["btn_help"][self.lang])
        self.lbl_delay.configure(text=self.texts["lbl_delay"][self.lang])
        self.lbl_max_pages.configure(text=self.texts["lbl_max_pages"][self.lang])
        self.lbl_max_depth.configure(text=self.texts["lbl_max_depth"][self.lang])
        self.lbl_format.configure(text=self.texts["lbl_format"][self.lang])
        self.lbl_concurrency.configure(text=self.texts["lbl_concurrency"][self.lang])
        self.cb_docs.configure(text=self.texts["cb_docs"][self.lang])
        self.cb_convert_docs.configure(text=self.texts["cb_convert_docs"][self.lang])
        self.lbl_mode.configure(text=self.texts["lbl_mode"][self.lang])
        self.lbl_folder.configure(text=self.texts["lbl_folder"][self.lang])
        self.btn_folder.configure(text=self.texts["btn_folder"][self.lang])

        self.cb_hybrid.configure(text=self.texts["cb_hybrid"][self.lang])
        self.cb_traf.configure(text=self.texts["cb_traf"][self.lang])
        self.cb_sitemap.configure(text=self.texts["cb_sitemap"][self.lang])
        self.cb_robots.configure(text=self.texts["cb_robots"][self.lang])
        self.cb_strict.configure(text=self.texts["cb_strict"][self.lang])
        self.lbl_exclude.configure(text=self.texts["lbl_exclude"][self.lang])
        self.lbl_require.configure(text=self.texts["lbl_require"][self.lang])

        self.cb_rm_email.configure(text=self.texts["cb_rm_email"][self.lang])
        self.cb_rm_phone.configure(text=self.texts["cb_rm_phone"][self.lang])
        self.cb_rm_pnr.configure(text=self.texts["cb_rm_pnr"][self.lang])
        self.cb_rm_ip.configure(text=self.texts["cb_rm_ip"][self.lang])

        self.start_btn.configure(text=self.texts["btn_start"][self.lang])
        if self.crawler_instance and self.crawler_instance.state == CrawlerState.PAUSED:
            self.pause_btn.configure(text=self.texts["btn_resume"][self.lang])
        else:
            self.pause_btn.configure(text=self.texts["btn_pause"][self.lang])
        self.stop_btn.configure(text=self.texts["btn_stop"][self.lang])

        self.tree.heading('Status', text=self.texts["col_status"][self.lang])
        self.tree.heading('Titel', text=self.texts["col_title"][self.lang])
        if not self.crawler_instance or self.crawler_instance.state in (
            CrawlerState.IDLE, CrawlerState.STOPPED
        ):
            self.stats_label.configure(text=self.texts["status_wait"][self.lang])

        self.headless_menu.configure(values=[v[self.lang] for v in self.texts["run_modes"].values()])
        self.headless_var.set(self.texts["run_modes"][internal_mode][self.lang])

        format_vals = [".json", ".md", ".txt", "Ingen text" if self.lang == "sv" else "No text"]
        self.format_menu.configure(values=format_vals)
        if self.format_var.get() not in format_vals:
            self.format_var.set(format_vals[-1])

    def _update_treeview_style(self, mode):
        style = ttk.Style()
        style.theme_use("default")
        if mode == "Dark":
            style.configure("Treeview", background="#2b2b2b", foreground="white",
                            fieldbackground="#2b2b2b", borderwidth=0, rowheight=25)
            style.configure("Treeview.Heading", background="#565b5e", foreground="white",
                            font=('Arial', 10, 'bold'), relief="flat")
            style.map('Treeview', background=[('selected', '#1f538d')])
            style.map("Treeview.Heading", background=[('active', '#343638')])
        else:
            style.configure("Treeview", background="#ffffff", foreground="black",
                            fieldbackground="#ffffff", borderwidth=0, rowheight=25)
            style.configure("Treeview.Heading", background="#e5e5e5", foreground="black",
                            font=('Arial', 10, 'bold'), relief="flat")
            style.map('Treeview', background=[('selected', '#3a7ebf')])
            style.map("Treeview.Heading", background=[('active', '#d1d1d1')])

    def open_help_window(self):
        help_win = ctk.CTkToplevel(self.root)
        help_win.title(self.texts["help_title"][self.lang])
        help_win.geometry("600x600")
        help_text = ctk.CTkTextbox(help_win, wrap=tk.WORD,
                                   font=ctk.CTkFont(family="Arial", size=13))
        help_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        help_text.insert(tk.END, self.texts["help_content"][self.lang])
        help_text.configure(state="disabled")

    def choose_directory(self):
        d = filedialog.askdirectory()
        if d:
            self.dir_entry.configure(state="normal")
            self.dir_var.set(d)
            self.dir_entry.configure(state="readonly")

    def _on_convert_docs_toggled(self):
        """Auto-aktivera dokumentnedladdning om konvertering kryssas i."""
        if self.convert_docs_var.get():
            self.docs_var.set(True)

    def _find_templates(self) -> Dict:
        if getattr(sys, 'frozen', False):
            script_dir = os.path.dirname(sys.executable)
        else:
            script_dir = os.path.dirname(os.path.abspath(__file__))
        json_path = os.path.join(script_dir, "sites.json")
        if os.path.isfile(json_path):
            try:
                with open(json_path, 'r', encoding='utf-8') as f:
                    sites = json.load(f)
                if isinstance(sites, list) and sites:
                    return {site.get("name", f"Sajt {i+1}"): site
                            for i, site in enumerate(sites)}
            except Exception:
                pass
        return {}

    def _apply_template(self, choice):
        none_label = self.texts["template_none"][self.lang]
        if choice == none_label:
            return

        site = self.templates.get(choice)
        if not site:
            return

        if site.get("start_url"):
            self.url_entry.delete(0, tk.END)
            self.url_entry.insert(0, site["start_url"])
        if "delay" in site:
            self.delay_entry.delete(0, tk.END)
            self.delay_entry.insert(0, str(site["delay"]))
        if "max_pages" in site:
            self.max_pages_entry.delete(0, tk.END)
            self.max_pages_entry.insert(0, str(site["max_pages"]))
        if "max_depth" in site:
            self.max_depth_entry.delete(0, tk.END)
            self.max_depth_entry.insert(0, str(site["max_depth"]))
        if "concurrency" in site:
            self.concurrency_entry.delete(0, tk.END)
            self.concurrency_entry.insert(0, str(site["concurrency"]))
        if "save_format" in site:
            self.format_var.set(site["save_format"])
        if "download_docs" in site:
            self.docs_var.set(site["download_docs"])
        if "headless_mode" in site:
            mode_key = site["headless_mode"]
            if mode_key in self.texts["run_modes"]:
                self.headless_var.set(self.texts["run_modes"][mode_key][self.lang])

        if "use_hybrid" in site:
            self.hybrid_var.set(site["use_hybrid"])
        if "use_trafilatura" in site:
            self.traf_var.set(site["use_trafilatura"])
        if "find_sitemap" in site:
            self.sitemap_var.set(site["find_sitemap"])
        if "respect_robots" in site:
            self.robots_var.set(site["respect_robots"])
        if "strict_domain" in site:
            self.strict_var.set(site["strict_domain"])

        if "exclude_keywords" in site:
            self.exclude_entry.delete(0, tk.END)
            kws = site["exclude_keywords"]
            self.exclude_entry.insert(0, ", ".join(kws) if isinstance(kws, list) else kws)
        if "require_keywords" in site:
            self.require_entry.delete(0, tk.END)
            kws = site["require_keywords"]
            self.require_entry.insert(0, ", ".join(kws) if isinstance(kws, list) else kws)

        if "remove_email" in site: self.rm_email_var.set(site["remove_email"])
        if "remove_phone" in site: self.rm_phone_var.set(site["remove_phone"])
        if "remove_pnr" in site: self.rm_pnr_var.set(site["remove_pnr"])
        if "remove_ip" in site: self.rm_ip_var.set(site["remove_ip"])
        if "convert_docs_to_md" in site:
            self.convert_docs_var.set(site["convert_docs_to_md"])
            if site["convert_docs_to_md"]:
                self.docs_var.set(True)

        if site.get("name"):
            base_dir = os.path.join(os.path.expanduser("~"), "Desktop", "crawl_output")
            safe_folder_name = slugify(site["name"])
            self.dir_var.set(os.path.join(base_dir, safe_folder_name))

        self._log_to_gui(self.texts["template_loaded"][self.lang].format(choice))

    def _log_to_gui(self, msg):
        t = datetime.now().strftime("%H:%M:%S")
        self.log_area.insert(tk.END, f"[{t}] {msg}\n")
        self.log_area.see(tk.END)

    def _build_ui(self):
        main_frame = ctk.CTkFrame(self.root, fg_color="transparent")
        main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)

        # Top URL bar
        url_frame = ctk.CTkFrame(main_frame)
        url_frame.pack(fill=tk.X, pady=(0, 10))
        self.lbl_url = ctk.CTkLabel(url_frame, text=self.texts["lbl_url"][self.lang],
                                    font=ctk.CTkFont(weight="bold"))
        self.lbl_url.pack(side=tk.LEFT, padx=(15, 10), pady=15)
        self.url_entry = ctk.CTkEntry(url_frame, width=250, placeholder_text="https://...")
        self.url_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 15), pady=15)
        self.url_entry.insert(0, "https://")

        self.lang_var = ctk.StringVar(value="🇸🇪 SV")
        self.lang_switch = ctk.CTkSegmentedButton(
            url_frame, values=["🇸🇪 SV", "🇬🇧 EN"],
            variable=self.lang_var, command=self.change_language_event
        )
        self.lang_switch.pack(side=tk.RIGHT, padx=15, pady=15)
        self.theme_switch = ctk.CTkSwitch(url_frame, text="☀️", width=40,
                                          command=self.change_appearance_mode_event)
        self.theme_switch.pack(side=tk.RIGHT, padx=(0, 15), pady=15)
        self.theme_switch.deselect()
        self.help_btn = ctk.CTkButton(
            url_frame, text=self.texts["btn_help"][self.lang], width=80,
            fg_color=("#d9d9d9", "#4a4a4a"),
            text_color=("black", "white"),
            hover_color=("#c9c9c9", "#5a5a5a"),
            command=self.open_help_window
        )
        self.help_btn.pack(side=tk.RIGHT, padx=(0, 15), pady=15)

        # Templates
        self.templates = self._find_templates()
        if self.templates:
            tpl_frame = ctk.CTkFrame(main_frame, fg_color="transparent")
            tpl_frame.pack(fill=tk.X, pady=(0, 5))
            self.lbl_template = ctk.CTkLabel(
                tpl_frame, text=self.texts["lbl_template"][self.lang],
                font=ctk.CTkFont(weight="bold"))
            self.lbl_template.pack(side=tk.LEFT, padx=(15, 10))
            none_label = self.texts["template_none"][self.lang]
            template_names = [none_label] + list(self.templates.keys())
            self.template_var = ctk.StringVar(value=none_label)
            self.template_menu = ctk.CTkOptionMenu(
                tpl_frame, variable=self.template_var, values=template_names,
                width=300, command=self._apply_template)
            self.template_menu.pack(side=tk.LEFT, padx=(0, 15))

        # Tabs
        self.tabview = ctk.CTkTabview(main_frame, height=260)
        self.tabview.pack(fill=tk.X, pady=(0, 10))
        tab_basic = self.tabview.add(self.texts["tab_basic"][self.lang])
        tab_adv = self.tabview.add(self.texts["tab_adv"][self.lang])

        # Tab 1: Basic
        tab_basic.grid_columnconfigure(1, weight=1)
        tab_basic.grid_columnconfigure(3, weight=1)

        self.lbl_delay = ctk.CTkLabel(tab_basic, text=self.texts["lbl_delay"][self.lang])
        self.lbl_delay.grid(row=0, column=0, padx=(10, 5), pady=8, sticky="e")
        self.delay_entry = ctk.CTkEntry(tab_basic, width=80)
        self.delay_entry.insert(0, "0.5")
        self.delay_entry.grid(row=0, column=1, padx=(0, 20), pady=8, sticky="w")

        self.lbl_max_pages = ctk.CTkLabel(tab_basic, text=self.texts["lbl_max_pages"][self.lang])
        self.lbl_max_pages.grid(row=0, column=2, padx=(10, 5), pady=8, sticky="e")
        self.max_pages_entry = ctk.CTkEntry(tab_basic, width=80)
        self.max_pages_entry.insert(0, "0")
        self.max_pages_entry.grid(row=0, column=3, padx=(0, 10), pady=8, sticky="w")

        self.lbl_max_depth = ctk.CTkLabel(tab_basic, text=self.texts["lbl_max_depth"][self.lang])
        self.lbl_max_depth.grid(row=1, column=0, padx=(10, 5), pady=8, sticky="e")
        self.max_depth_entry = ctk.CTkEntry(tab_basic, width=80)
        self.max_depth_entry.insert(0, "0")
        self.max_depth_entry.grid(row=1, column=1, padx=(0, 20), pady=8, sticky="w")

        self.lbl_format = ctk.CTkLabel(tab_basic, text=self.texts["lbl_format"][self.lang])
        self.lbl_format.grid(row=1, column=2, padx=(10, 5), pady=8, sticky="e")
        # Default: .md (LLM-vänligt, fungerar bäst med RAG-pipelines som
        # chunkar i stycken). .json finns kvar för dem som vill ha strukturerade
        # fält per sida.
        self.format_var = ctk.StringVar(value=".md")
        self.format_menu = ctk.CTkOptionMenu(
            tab_basic, variable=self.format_var,
            values=[".json", ".md", ".txt", "Ingen text"], width=100)
        self.format_menu.grid(row=1, column=3, padx=(0, 10), pady=8, sticky="w")

        self.lbl_concurrency = ctk.CTkLabel(tab_basic, text=self.texts["lbl_concurrency"][self.lang])
        self.lbl_concurrency.grid(row=2, column=0, padx=(10, 5), pady=8, sticky="e")
        self.concurrency_entry = ctk.CTkEntry(tab_basic, width=80)
        self.concurrency_entry.insert(0, "10")
        self.concurrency_entry.grid(row=2, column=1, padx=(0, 20), pady=8, sticky="w")

        # Default: ladda ner och konvertera dokument — det är den vanligaste
        # användningen för RAG-system. Användaren kan stänga av om de bara
        # vill ha sidor.
        self.docs_var = ctk.BooleanVar(value=True)
        self.cb_docs = ctk.CTkCheckBox(tab_basic, text=self.texts["cb_docs"][self.lang],
                                       variable=self.docs_var)
        self.cb_docs.grid(row=2, column=2, padx=10, pady=8, sticky="w")

        self.convert_docs_var = ctk.BooleanVar(value=True)
        self.cb_convert_docs = ctk.CTkCheckBox(
            tab_basic, text=self.texts["cb_convert_docs"][self.lang],
            variable=self.convert_docs_var,
            command=self._on_convert_docs_toggled)
        self.cb_convert_docs.grid(row=2, column=3, padx=10, pady=8, sticky="w")

        self.lbl_mode = ctk.CTkLabel(tab_basic, text=self.texts["lbl_mode"][self.lang])
        self.lbl_mode.grid(row=3, column=0, padx=(10, 5), pady=8, sticky="e")
        self.headless_var = ctk.StringVar(value=self.texts["run_modes"]["headless"][self.lang])
        self.headless_menu = ctk.CTkOptionMenu(
            tab_basic, variable=self.headless_var,
            values=[v[self.lang] for v in self.texts["run_modes"].values()], width=180)
        self.headless_menu.grid(row=3, column=1, columnspan=3, padx=(0, 10), pady=8, sticky="w")

        self.lbl_folder = ctk.CTkLabel(tab_basic, text=self.texts["lbl_folder"][self.lang])
        self.lbl_folder.grid(row=4, column=0, padx=(10, 5), pady=8, sticky="e")
        self.dir_var = ctk.StringVar(
            value=os.path.join(os.path.expanduser("~"), "Desktop", "crawl_output"))
        self.dir_entry = ctk.CTkEntry(tab_basic, textvariable=self.dir_var, state="readonly")
        self.dir_entry.grid(row=4, column=1, columnspan=2, sticky="ew", padx=(0, 10), pady=8)
        self.btn_folder = ctk.CTkButton(
            tab_basic, text=self.texts["btn_folder"][self.lang],
            width=100, command=self.choose_directory)
        self.btn_folder.grid(row=4, column=3, padx=(0, 10), pady=8, sticky="w")

        # Tab 2: Advanced
        tab_adv.grid_columnconfigure(1, weight=1)

        self.hybrid_var = ctk.BooleanVar(value=True)
        self.cb_hybrid = ctk.CTkCheckBox(tab_adv, text=self.texts["cb_hybrid"][self.lang],
                                         variable=self.hybrid_var)
        self.cb_hybrid.grid(row=0, column=0, padx=10, pady=5, sticky="w")

        self.traf_var = ctk.BooleanVar(value=HAS_TRAFILATURA)
        self.cb_traf = ctk.CTkCheckBox(tab_adv, text=self.texts["cb_traf"][self.lang],
                                       variable=self.traf_var)
        self.cb_traf.grid(row=0, column=1, padx=10, pady=5, sticky="w")
        if not HAS_TRAFILATURA:
            self.cb_traf.configure(state="disabled")

        self.sitemap_var = ctk.BooleanVar(value=True)
        self.cb_sitemap = ctk.CTkCheckBox(tab_adv, text=self.texts["cb_sitemap"][self.lang],
                                          variable=self.sitemap_var)
        self.cb_sitemap.grid(row=1, column=0, padx=10, pady=5, sticky="w")

        self.robots_var = ctk.BooleanVar(value=True)
        self.cb_robots = ctk.CTkCheckBox(tab_adv, text=self.texts["cb_robots"][self.lang],
                                         variable=self.robots_var)
        self.cb_robots.grid(row=1, column=1, padx=10, pady=5, sticky="w")

        self.strict_var = ctk.BooleanVar(value=True)
        self.cb_strict = ctk.CTkCheckBox(tab_adv, text=self.texts["cb_strict"][self.lang],
                                         variable=self.strict_var)
        self.cb_strict.grid(row=1, column=2, padx=10, pady=5, sticky="w")

        self.lbl_exclude = ctk.CTkLabel(tab_adv, text=self.texts["lbl_exclude"][self.lang])
        self.lbl_exclude.grid(row=2, column=0, padx=(10, 5), pady=5, sticky="e")
        self.exclude_entry = ctk.CTkEntry(tab_adv, placeholder_text="images, login, kalender")
        self.exclude_entry.grid(row=2, column=1, columnspan=2, sticky="ew", padx=(0, 10), pady=5)

        self.lbl_require = ctk.CTkLabel(tab_adv, text=self.texts["lbl_require"][self.lang])
        self.lbl_require.grid(row=3, column=0, padx=(10, 5), pady=5, sticky="e")
        self.require_entry = ctk.CTkEntry(tab_adv, placeholder_text="intranat, bibliotek")
        self.require_entry.grid(row=3, column=1, columnspan=2, sticky="ew", padx=(0, 10), pady=5)

        arow4 = ctk.CTkFrame(tab_adv, fg_color="transparent")
        arow4.grid(row=4, column=0, columnspan=3, pady=(10, 0), sticky="w")

        self.rm_email_var = ctk.BooleanVar(value=False)
        self.cb_rm_email = ctk.CTkCheckBox(arow4, text=self.texts["cb_rm_email"][self.lang],
                                           variable=self.rm_email_var)
        self.cb_rm_email.grid(row=0, column=0, padx=10, pady=5, sticky="w")

        self.rm_phone_var = ctk.BooleanVar(value=False)
        self.cb_rm_phone = ctk.CTkCheckBox(arow4, text=self.texts["cb_rm_phone"][self.lang],
                                           variable=self.rm_phone_var)
        self.cb_rm_phone.grid(row=0, column=1, padx=10, pady=5, sticky="w")

        self.rm_pnr_var = ctk.BooleanVar(value=False)
        self.cb_rm_pnr = ctk.CTkCheckBox(arow4, text=self.texts["cb_rm_pnr"][self.lang],
                                         variable=self.rm_pnr_var)
        self.cb_rm_pnr.grid(row=0, column=2, padx=10, pady=5, sticky="w")

        self.rm_ip_var = ctk.BooleanVar(value=False)
        self.cb_rm_ip = ctk.CTkCheckBox(arow4, text=self.texts["cb_rm_ip"][self.lang],
                                        variable=self.rm_ip_var)
        self.cb_rm_ip.grid(row=1, column=0, padx=10, pady=5, sticky="w")

        # Buttons
        btn_frame = ctk.CTkFrame(main_frame, fg_color="transparent")
        btn_frame.pack(pady=5)
        self.start_btn = ctk.CTkButton(
            btn_frame, text=self.texts["btn_start"][self.lang],
            font=ctk.CTkFont(weight="bold"),
            fg_color="#1f6aa5", command=self.start_crawl)
        self.start_btn.pack(side=tk.LEFT, padx=10)
        self.pause_btn = ctk.CTkButton(
            btn_frame, text=self.texts["btn_pause"][self.lang],
            state="disabled",
            fg_color=("#d9d9d9", "#4a4a4a"),
            text_color=("black", "white"),
            hover_color=("#c9c9c9", "#5a5a5a"),
            command=self.toggle_pause)
        self.pause_btn.pack(side=tk.LEFT, padx=10)
        self.stop_btn = ctk.CTkButton(
            btn_frame, text=self.texts["btn_stop"][self.lang],
            state="disabled",
            fg_color=("#d35b5b", "#a51f1f"),
            hover_color=("#c42b2b", "#8a1a1a"),
            text_color=("white", "white"),
            command=self.stop_crawl)
        self.stop_btn.pack(side=tk.LEFT, padx=10)

        # Stats & display
        self.stats_label = ctk.CTkLabel(
            main_frame, text=self.texts["status_wait"][self.lang],
            font=ctk.CTkFont(family="Consolas", size=12, weight="bold"),
            text_color="#4caf50")
        self.stats_label.pack(fill=tk.X, pady=5)
        self.progress_bar = ctk.CTkProgressBar(main_frame, orientation="horizontal")

        table_frame = ctk.CTkFrame(main_frame)
        table_frame.pack(fill=tk.BOTH, expand=True)
        self.tree = ttk.Treeview(table_frame, columns=('URL', 'Status', 'Titel'),
                                 show='headings', height=7)
        self.tree.heading('URL', text='URL')
        self.tree.heading('Status', text=self.texts["col_status"][self.lang])
        self.tree.heading('Titel', text=self.texts["col_title"][self.lang])
        self.tree.column('URL', width=300)
        self.tree.column('Status', width=130)
        self.tree.column('Titel', width=300)

        scrollbar = ctk.CTkScrollbar(table_frame, orientation="vertical",
                                     command=self.tree.yview)
        self.tree.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.tree.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        self.tree.bind(
            "<Double-1>",
            lambda e: webbrowser.open(self.tree.item(self.tree.selection()[0])['values'][0])
            if self.tree.selection()
            and self.tree.item(self.tree.selection()[0])['values'][0].startswith("http")
            else None
        )

        self.log_area = ctk.CTkTextbox(main_frame,
                                       font=ctk.CTkFont(family="Consolas", size=11),
                                       height=100)
        self.log_area.pack(fill=tk.BOTH, expand=False, pady=(10, 0))

    def process_queue(self):
        try:
            for _ in range(20):
                try:
                    msg_type, data = self.msg_queue.get_nowait()
                    if msg_type == "log":
                        t = datetime.now().strftime("%H:%M:%S")
                        self.log_area.insert(tk.END, f"[{t}] {data}\n")
                        self.log_area.see(tk.END)
                        if int(self.log_area.index('end-1c').split('.')[0]) > 500:
                            self.log_area.delete("1.0", "2.0")
                    elif msg_type == "table":
                        self.tree.insert('', 0, values=data)
                        if len(self.tree.get_children()) > 100:
                            self.tree.delete(self.tree.get_children()[-1])
                    elif msg_type == "stats_data":
                        self.stats_label.configure(
                            text=self.texts["stats_fmt"][self.lang].format(*data))
                        if self.progress_bar and self.max_pages_entry.get() != "0":
                            try:
                                max_p = int(self.max_pages_entry.get())
                                if max_p > 0:
                                    self.progress_bar.set(data[0] / max_p)
                            except ValueError:
                                pass
                    elif msg_type == "login_wait":
                        messagebox.showinfo(
                            "Inloggning / Login",
                            "Logga in i webbläsaren. Tryck OK här när du är klar!"
                            if self.lang == "sv"
                            else "Please login in the browser. Click OK here when done!"
                        )
                        if self.crawler_instance:
                            self.crawler_instance.login_event.set()
                    elif msg_type == "done":
                        self.start_btn.configure(state="normal")
                        self.pause_btn.configure(state="disabled")
                        self.stop_btn.configure(state="disabled")
                        self.progress_bar.stop()
                        self.progress_bar.pack_forget()
                except queue.Empty:
                    break
                except Exception as inner_e:
                    self.log_area.insert(tk.END, f"[GUI FEL] Kunde inte rita rad: {inner_e}\n")
        finally:
            self.root.after(100, self.process_queue)

    def start_crawl(self):
        url = self.url_entry.get().strip()
        if not url or url == "https://":
            return
        try:
            delay_val = float(self.delay_entry.get())
            max_pages_val = int(self.max_pages_entry.get())
            max_depth_val = int(self.max_depth_entry.get())
            concurrency_val = int(self.concurrency_entry.get())
        except ValueError:
            messagebox.showerror(
                "Ogiltigt värde" if self.lang == "sv" else "Invalid value",
                ("Kontrollera att Fördröjning, Max sidor, Max djup och "
                 "Samtidighet är giltiga tal (använd punkt som decimalavgränsare).")
                if self.lang == "sv" else
                ("Please check that Delay, Max pages, Max depth and "
                 "Concurrency are valid numbers (use dot as decimal separator).")
            )
            return

        self.start_btn.configure(state="disabled")
        self.pause_btn.configure(state="normal", text=self.texts["btn_pause"][self.lang])
        self.stop_btn.configure(state="normal")
        self.tree.delete(*self.tree.get_children())
        self.log_area.delete("1.0", tk.END)

        self.progress_bar.pack(fill=tk.X, pady=(0, 10))
        if max_pages_val > 0:
            self.progress_bar.configure(mode='determinate')
            self.progress_bar.set(0)
        else:
            self.progress_bar.configure(mode='indeterminate')
            self.progress_bar.start()

        inverted_map = {v[self.lang]: k for k, v in self.texts["run_modes"].items()}

        config = {
            "start_url": url,
            "output_dir": self.dir_var.get(),
            "delay": delay_val,
            "max_pages": max_pages_val,
            "max_depth": max_depth_val,
            "concurrency": concurrency_val,
            "save_format": self.format_var.get(),
            "headless_mode": inverted_map.get(self.headless_var.get(), "headless"),
            "respect_robots": self.robots_var.get(),
            "find_sitemap": self.sitemap_var.get(),
            "use_hybrid": self.hybrid_var.get(),
            "use_trafilatura": self.traf_var.get(),
            "download_docs": self.docs_var.get(),
            "strict_domain": self.strict_var.get(),
            "exclude_keywords": [k.strip().lower()
                                 for k in self.exclude_entry.get().split(",") if k.strip()],
            "require_keywords": [k.strip().lower()
                                 for k in self.require_entry.get().split(",") if k.strip()],
            "remove_email": self.rm_email_var.get(),
            "remove_phone": self.rm_phone_var.get(),
            "remove_pnr": self.rm_pnr_var.get(),
            "remove_ip": self.rm_ip_var.get(),
            "convert_docs_to_md": self.convert_docs_var.get(),
            "incremental": True,
        }

        self.crawler_instance = AsyncWebCrawler(config, self.msg_queue)
        threading.Thread(
            target=lambda: asyncio.run(self.crawler_instance.crawl()),
            daemon=True
        ).start()

    def toggle_pause(self):
        if self.crawler_instance:
            is_paused = self.crawler_instance.pause()
            self.pause_btn.configure(
                text=self.texts["btn_resume"][self.lang]
                if is_paused else self.texts["btn_pause"][self.lang]
            )

    def stop_crawl(self):
        if self.crawler_instance:
            self.stop_btn.configure(state="disabled")
            self.crawler_instance.stop()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", type=str)
    parser.add_argument("--webhook", type=str)
    args, _ = parser.parse_known_args()

    if HAS_UVLOOP:
        uvloop.install()

    if args.config:
        asyncio.run(run_cli_mode(args.config,
                                 args.webhook or os.environ.get("WEBHOOK_URL")))
    else:
        AppGUI(ctk.CTk()).root.mainloop()


if __name__ == "__main__":
    main()